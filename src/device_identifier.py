#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
设备识别模块 - 任务5：创建设备识别模块
从PPT中识别设备组和对应的pdid标签
"""

from pptx import Presentation
import pandas as pd
import json
from typing import Dict, List, Tuple, Optional

class DeviceIdentifier:
    """设备识别器"""
    
    def __init__(self, ppt_path: str, excel_path: str = None):
        """
        初始化设备识别器
        
        Args:
            ppt_path: PPT文件路径
            excel_path: Excel文件路径（可选，用于验证）
        """
        self.ppt_path = ppt_path
        self.excel_path = excel_path
        self.presentation = None
        self.excel_data = None
        
    def load_presentation(self) -> bool:
        """
        加载PPT文件
        
        Returns:
            bool: 是否成功加载
        """
        try:
            self.presentation = Presentation(self.ppt_path)
            print(f"✅ 成功加载PPT文件: {self.ppt_path}")
            print(f"📊 幻灯片数量: {len(self.presentation.slides)}")
            return True
        except Exception as e:
            print(f"❌ 加载PPT文件失败: {e}")
            return False
    
    def load_excel_data(self) -> bool:
        """
        加载Excel数据（用于验证）
        
        Returns:
            bool: 是否成功加载
        """
        if not self.excel_path:
            print("⚠️ 未提供Excel文件路径，跳过数据验证")
            return True
        
        try:
            self.excel_data = pd.read_excel(self.excel_path)
            print(f"✅ 成功加载Excel文件: {self.excel_path}")
            print(f"📊 数据形状: {self.excel_data.shape}")
            return True
        except Exception as e:
            print(f"❌ 加载Excel文件失败: {e}")
            return False
    
    def identify_pdid_labels(self) -> Dict[int, List[Dict]]:
        """
        识别PPT中的pdid标签
        
        Returns:
            Dict[int, List[Dict]]: 幻灯片索引到pdid标签信息的映射
        """
        if self.presentation is None:
            return {}
        
        pdid_labels = {}
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            print(f"\n🔍 识别第{slide_idx + 1}张幻灯片中的pdid标签:")
            
            slide_labels = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    
                    # 检查是否为pdid标签
                    if text.startswith('pdid:'):
                        try:
                            pdid = int(text.split(':')[1].strip())
                            label_info = {
                                'shape': shape,
                                'name': shape.name,
                                'text': text,
                                'pdid': pdid,
                                'position': {
                                    'left': shape.left,
                                    'top': shape.top,
                                    'width': shape.width,
                                    'height': shape.height
                                }
                            }
                            slide_labels.append(label_info)
                            print(f"   ✅ 发现pdid标签: {text} (形状: {shape.name})")
                        except (ValueError, IndexError):
                            print(f"   ⚠️ 无法解析pdid标签: {text}")
            
            pdid_labels[slide_idx] = slide_labels
            print(f"   📊 本页pdid标签数量: {len(slide_labels)}")
        
        return pdid_labels
    
    def identify_device_groups(self) -> Dict[int, List[Dict]]:
        """
        识别PPT中的设备组
        
        Returns:
            Dict[int, List[Dict]]: 幻灯片索引到设备组信息的映射
        """
        if self.presentation is None:
            return {}
        
        device_groups = {}
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            print(f"\n🔍 识别第{slide_idx + 1}张幻灯片中的设备组:")
            
            slide_groups = []
            for shape in slide.shapes:
                shape_info = {
                    'shape': shape,
                    'name': shape.name,
                    'type': type(shape).__name__,
                    'has_text': shape.has_text_frame,
                    'text': shape.text if shape.has_text_frame else "",
                    'position': {
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    }
                }
                
                # 判断是否为设备组相关形状
                is_device_group = False
                
                # 根据形状名称判断
                if 'smart_home_switch' in shape.name.lower() or 'switch' in shape.name.lower():
                    is_device_group = True
                
                # 根据文本内容判断
                elif shape.has_text_frame:
                    # 安全地获取文本内容
                    try:
                        text_content = shape.text if hasattr(shape, 'text') else ""
                        if '开关' in text_content or 'switch' in text_content.lower():
                            is_device_group = True
                    except:
                        text_content = ""
                
                if is_device_group:
                    slide_groups.append(shape_info)
                    # 安全地显示文本内容
                    try:
                        display_text = shape.text[:30] if hasattr(shape, 'text') else ""
                        print(f"   ✅ 发现设备组: {shape.name} - {display_text}...")
                    except:
                        print(f"   ✅ 发现设备组: {shape.name}")
            
            device_groups[slide_idx] = slide_groups
            print(f"   📊 本页设备组数量: {len(slide_groups)}")
        
        return device_groups
    
    def match_devices_with_pdid(self, device_groups: Dict, pdid_labels: Dict) -> Dict[int, List[Dict]]:
        """
        将设备组与pdid标签进行匹配
        
        Args:
            device_groups: 设备组信息
            pdid_labels: pdid标签信息
            
        Returns:
            Dict[int, List[Dict]]: 匹配结果
        """
        matched_devices = {}
        
        for slide_idx in device_groups.keys():
            print(f"\n🎯 匹配第{slide_idx + 1}张幻灯片中的设备组和pdid标签:")
            
            slide_devices = device_groups.get(slide_idx, [])
            slide_labels = pdid_labels.get(slide_idx, [])
            
            matched = []
            
            for device in slide_devices:
                device_pos = device['position']
                
                # 查找与设备组位置相近的pdid标签
                matched_pdid = None
                for label in slide_labels:
                    label_pos = label['position']
                    
                    # 检查pdid标签是否在设备组下方
                    if (label_pos['top'] >= device_pos['top'] + device_pos['height'] and
                        label_pos['left'] >= device_pos['left'] and
                        label_pos['left'] + label_pos['width'] <= device_pos['left'] + device_pos['width']):
                        
                        matched_pdid = label['pdid']
                        print(f"   ✅ 设备组 {device['name']} 匹配pdid: {matched_pdid}")
                        break
                
                if matched_pdid:
                    device['matched_pdid'] = matched_pdid
                    matched.append(device)
                else:
                    print(f"   ⚠️ 设备组 {device['name']} 未找到匹配的pdid标签")
            
            matched_devices[slide_idx] = matched
            print(f"   📊 本页匹配成功设备组数量: {len(matched)}")
        
        return matched_devices
    
    def validate_with_excel(self, matched_devices: Dict) -> Dict[int, List[Dict]]:
        """
        使用Excel数据验证匹配结果
        
        Args:
            matched_devices: 匹配的设备组信息
            
        Returns:
            Dict[int, List[Dict]]: 验证结果
        """
        if self.excel_data is None:
            print("⚠️ 未提供Excel数据，跳过验证")
            return matched_devices
        
        validated_devices = {}
        
        for slide_idx, devices in matched_devices.items():
            print(f"\n🔍 验证第{slide_idx + 1}张幻灯片中的设备组:")
            
            validated = []
            for device in devices:
                pdid = device.get('matched_pdid')
                
                if pdid:
                    # 在Excel中查找对应的产品信息
                    product_info = self.excel_data[self.excel_data['产品ID'] == pdid]
                    
                    if not product_info.empty:
                        device['excel_validation'] = {
                            'valid': True,
                            'device_name': product_info['设备名称'].iloc[0],
                            'brand': product_info['品牌'].iloc[0],
                            'spec': product_info['主规格'].iloc[0] if '主规格' in product_info.columns else ''
                        }
                        print(f"   ✅ 设备组 {device['name']} (pdid: {pdid}) 验证成功")
                    else:
                        device['excel_validation'] = {
                            'valid': False,
                            'error': f"Excel中未找到产品ID {pdid}"
                        }
                        print(f"   ❌ 设备组 {device['name']} (pdid: {pdid}) 验证失败")
                
                validated.append(device)
            
            validated_devices[slide_idx] = validated
        
        return validated_devices
    
    def generate_identification_report(self, matched_devices: Dict) -> Dict:
        """
        生成设备识别报告
        
        Args:
            matched_devices: 匹配的设备组信息
            
        Returns:
            Dict: 识别报告
        """
        report = {
            'total_slides': len(self.presentation.slides),
            'total_devices_identified': 0,
            'total_pdid_labels_found': 0,
            'successful_matches': 0,
            'failed_matches': 0,
            'slide_details': {},
            'summary': {}
        }
        
        for slide_idx, devices in matched_devices.items():
            slide_report = {
                'slide_number': slide_idx + 1,
                'devices_count': len(devices),
                'devices': []
            }
            
            for device in devices:
                device_report = {
                    'shape_name': device['name'],
                    'device_text': device['text'][:50],
                    'matched_pdid': device.get('matched_pdid'),
                    'position': device['position']
                }
                
                if 'excel_validation' in device:
                    device_report['excel_validation'] = device['excel_validation']
                
                slide_report['devices'].append(device_report)
                
                if device.get('matched_pdid'):
                    report['successful_matches'] += 1
                else:
                    report['failed_matches'] += 1
            
            report['total_devices_identified'] += len(devices)
            report['slide_details'][slide_idx] = slide_report
        
        # 统计pdid标签总数
        for slide_idx in range(len(self.presentation.slides)):
            slide = self.presentation.slides[slide_idx]
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip().startswith('pdid:'):
                    report['total_pdid_labels_found'] += 1
        
        # 生成摘要
        report['summary'] = {
            'identification_rate': f"{report['successful_matches'] / report['total_devices_identified'] * 100:.1f}%" if report['total_devices_identified'] > 0 else "0%",
            'average_devices_per_slide': report['total_devices_identified'] / len(self.presentation.slides) if len(self.presentation.slides) > 0 else 0
        }
        
        return report
    
    def identify_devices(self) -> Optional[Dict]:
        """
        执行设备识别流程
        
        Returns:
            Optional[Dict]: 识别报告
        """
        print("=" * 60)
        print("🔧 开始设备识别 - 任务5")
        print("=" * 60)
        
        # 加载PPT文件
        if not self.load_presentation():
            return None
        
        # 加载Excel数据（可选）
        if not self.load_excel_data():
            return None
        
        # 识别pdid标签
        pdid_labels = self.identify_pdid_labels()
        
        # 识别设备组
        device_groups = self.identify_device_groups()
        
        # 匹配设备组和pdid标签
        matched_devices = self.match_devices_with_pdid(device_groups, pdid_labels)
        
        # 使用Excel数据验证
        validated_devices = self.validate_with_excel(matched_devices)
        
        # 生成识别报告
        report = self.generate_identification_report(validated_devices)
        
        print("=" * 60)
        print("📊 设备识别报告摘要:")
        print(f"   总幻灯片数: {report['total_slides']}")
        print(f"   识别的设备组总数: {report['total_devices_identified']}")
        print(f"   发现的pdid标签总数: {report['total_pdid_labels_found']}")
        print(f"   成功匹配的设备组: {report['successful_matches']}")
        print(f"   匹配失败设备组: {report['failed_matches']}")
        print(f"   识别率: {report['summary']['identification_rate']}")
        print("=" * 60)
        
        return report

def identify_devices_in_ppt(ppt_path: str, excel_path: str = None) -> Optional[Dict]:
    """
    设备识别主函数
    
    Args:
        ppt_path: PPT文件路径
        excel_path: Excel文件路径（可选）
        
    Returns:
        Optional[Dict]: 识别报告
    """
    identifier = DeviceIdentifier(ppt_path, excel_path)
    return identifier.identify_devices()

if __name__ == "__main__":
    # 测试函数
    ppt_path = "E:\\Programs\\smarthome\\智能家居模具库.pptx"
    excel_path = "E:\\Programs\\smarthome\\智能家居模具库.xlsx"
    
    report = identify_devices_in_ppt(ppt_path, excel_path)
    
    if report:
        print("🎯 设备识别任务完成")
        
        # 保存报告到文件
        import json
        with open("E:\\Programs\\smarthome\\src\\device_identification_report.json", 'w', encoding='utf-8') as f:
            json.dump(report, f, ensure_ascii=False, indent=2)
        print("📄 识别报告已保存到: device_identification_report.json")
    else:
        print("❌ 设备识别失败")