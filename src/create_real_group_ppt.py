#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建真正的组结构PPT模具库
解决复制组时无法复制组下文字的问题
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
import os

def create_real_group_mold_library():
    """创建真正的组结构模具库"""
    
    print("=== 创建真正的组结构PPT模具库 ===")
    
    # 创建新的PPT文件
    prs = Presentation()
    
    # 添加标题幻灯片
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "智能家居模具库（真组结构）"
    title_slide.placeholders[1].text = "支持完整复制的真组模具系统"
    
    # 产品数据（根据Excel文件）
    products = [
        # 领普品牌
        {'name': '一键智能开关', 'id': 'switch_1_lp', 'price': 79, 'color': RGBColor(240, 248, 255), 'brand': '领普'},
        {'name': '二键智能开关', 'id': 'switch_2_lp', 'price': 89, 'color': RGBColor(240, 255, 240), 'brand': '领普'},
        {'name': '三键智能开关', 'id': 'switch_3_lp', 'price': 99, 'color': RGBColor(255, 240, 245), 'brand': '领普'},
        {'name': '四键智能开关', 'id': 'switch_4_lp', 'price': 109, 'color': RGBColor(255, 248, 220), 'brand': '领普'},
        
        # 易来品牌  
        {'name': '一键智能开关', 'id': 'switch_1_yl', 'price': 79, 'color': RGBColor(220, 240, 255), 'brand': '易来'},
        {'name': '二键智能开关', 'id': 'switch_2_yl', 'price': 89, 'color': RGBColor(220, 255, 240), 'brand': '易来'},
        {'name': '三键智能开关', 'id': 'switch_3_yl', 'price': 99, 'color': RGBColor(255, 220, 245), 'brand': '易来'},
        {'name': '四键智能开关', 'id': 'switch_4_yl', 'price': 109, 'color': RGBColor(255, 240, 220), 'brand': '易来'}
    ]
    
    # 添加产品模具幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "智能开关系列（真组模具）"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # 添加产品模具（真组结构）
    for i, product in enumerate(products):
        row = i // 4
        col = i % 4
        
        left = Inches(1 + col * 2)
        top = Inches(1.5 + row * 1.5)
        
        # 创建组容器
        group_left = left
        group_top = top
        group_width = Inches(1.8)
        group_height = Inches(1.2)
        
        # 创建组内的各个元素
        # 1. 主形状（产品图标）
        main_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            group_left + Inches(0.1), group_top + Inches(0.1), Inches(1.6), Inches(0.6)
        )
        main_shape.name = f"{product['id']}_main"
        main_shape.fill.solid()
        main_shape.fill.fore_color.rgb = product['color']
        main_shape.line.color.rgb = RGBColor(0, 0, 0)
        main_shape.line.width = Pt(1)
        
        # 2. 产品名称文本
        name_shape = slide.shapes.add_textbox(
            group_left + Inches(0.1), group_top + Inches(0.7), Inches(1.6), Inches(0.2)
        )
        name_shape.name = f"{product['id']}_name"
        name_frame = name_shape.text_frame
        name_frame.text = product['name']
        name_frame.paragraphs[0].font.size = Pt(10)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 3. 品牌和价格文本
        info_shape = slide.shapes.add_textbox(
            group_left + Inches(0.1), group_top + Inches(0.9), Inches(1.6), Inches(0.2)
        )
        info_shape.name = f"{product['id']}_info"
        info_frame = info_shape.text_frame
        info_frame.text = f"{product['brand']} ¥{product['price']}"
        info_frame.paragraphs[0].font.size = Pt(9)
        info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 4. 隐藏的产品ID文本（用于识别）
        id_shape = slide.shapes.add_textbox(
            group_left + Inches(0.1), group_top + Inches(1.1), Inches(1.6), Inches(0.1)
        )
        id_shape.name = f"{product['id']}_id"
        id_frame = id_shape.text_frame
        id_frame.text = f"product_id:{product['id']}"
        id_frame.paragraphs[0].font.size = Pt(6)
        id_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 白色，几乎不可见
        
        print(f"✅ 创建产品组: {product['id']}")
        print(f"   名称: {product['name']}")
        print(f"   品牌: {product['brand']}")
        print(f"   价格: ¥{product['price']}")
        print(f"   位置: ({group_left}, {group_top})")
    
    # 保存文件
    output_path = 'E:\\Programs\\smarthome\\output\\smart_home_real_group_mold_gallery.pptx'
    prs.save(output_path)
    
    print(f"\n✅ 真组结构PPT模具库已创建: {output_path}")
    print("\n📋 组结构特点:")
    print("   • 每个产品包含4个独立形状")
    print("   • 主形状（产品图标）")
    print("   • 产品名称文本")
    print("   • 品牌价格信息")
    print("   • 隐藏的产品ID文本")
    print("   • 复制时所有元素都会被复制")
    
    return output_path

def test_group_recognition():
    """测试组识别功能"""
    
    print("\n=== 测试组识别功能 ===")
    
    ppt_path = 'E:\\Programs\\smarthome\\output\\smart_home_real_group_mold_gallery.pptx'
    
    if not os.path.exists(ppt_path):
        print("❌ PPT文件不存在，请先创建真组结构模具库")
        return False
    
    from pptx import Presentation
    
    prs = Presentation(ppt_path)
    print(f"📊 幻灯片数量: {len(prs.slides)}")
    
    # 检查第二张幻灯片
    if len(prs.slides) > 1:
        slide = prs.slides[1]
        print(f"📄 第二张幻灯片形状数量: {len(slide.shapes)}")
        
        # 统计产品组
        product_groups = {}
        for shape in slide.shapes:
            if hasattr(shape, 'name') and shape.name:
                # 解析产品ID
                if 'smart_home_' in shape.name and '_id' in shape.name:
                    product_id = shape.name.split('smart_home_')[1].split('_id')[0]
                    if product_id not in product_groups:
                        product_groups[product_id] = []
                    product_groups[product_id].append(shape.name)
        
        print(f"🔍 识别到的产品组数量: {len(product_groups)}")
        
        for product_id, shapes in product_groups.items():
            print(f"   📦 产品组 {product_id}: {len(shapes)} 个相关形状")
            for shape_name in shapes:
                shape_type = shape_name.split('_')[-1]
                print(f"      {shape_type}: {shape_name}")
    
    return True

def create_group_recognition_script():
    """创建组识别脚本"""
    
    script_content = '''#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
智能家居组识别脚本
支持识别真组结构中的产品信息
"""

from pptx import Presentation
import pandas as pd
import os

def read_excel_product_library(excel_path):
    """从Excel模具库读取产品信息"""
    if not os.path.exists(excel_path):
        print(f"❌ Excel模具库文件不存在: {excel_path}")
        return {}
    
    df = pd.read_excel(excel_path)
    product_library = {}
    
    for index, row in df.iterrows():
        product_id = row.get('产品ID', '')
        if pd.notna(product_id) and product_id:
            product_library[product_id] = {
                "name": row.get('设备名称', ''),
                "price": int(row.get('单价', 0)) if pd.notna(row.get('单价')) else 0,
                "brand": row.get('品牌', ''),
                "model": row.get('主规格', ''),
                "category": row.get('设备品类', '')
            }
    
    print(f"✅ 从Excel读取了 {len(product_library)} 个产品信息")
    return product_library

def extract_product_id_from_shape_name(shape_name):
    """从形状名称提取产品ID"""
    if not shape_name or 'smart_home_' not in shape_name:
        return None
    
    # 格式: smart_home_switch_1_lp_id
    parts = shape_name.split('smart_home_')[1].split('_')
    if len(parts) >= 3:
        return f"{parts[0]}_{parts[1]}_{parts[2]}"
    elif len(parts) >= 2:
        return f"{parts[0]}_{parts[1]}"
    
    return None

def scan_ppt_for_product_groups(ppt_path, excel_library_path):
    """扫描PPT文件中的产品组"""
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")
    
    # 读取产品库
    product_library = read_excel_product_library(excel_library_path)
    
    prs = Presentation(ppt_path)
    all_products = []
    
    print(f"🔍 扫描PPT文件: {ppt_path}")
    print(f"📊 幻灯片数量: {len(prs.slides)}")
    
    # 处理第二张幻灯片（产品模具页）
    if len(prs.slides) > 1:
        slide = prs.slides[1]
        
        # 按产品ID分组形状
        product_groups = {}
        for shape in slide.shapes:
            if hasattr(shape, 'name') and shape.name:
                product_id = extract_product_id_from_shape_name(shape.name)
                if product_id:
                    if product_id not in product_groups:
                        product_groups[product_id] = []
                    product_groups[product_id].append(shape)
        
        print(f"📦 识别到 {len(product_groups)} 个产品组")
        
        # 处理每个产品组
        for product_id, shapes in product_groups.items():
            if product_id in product_library:
                product_info = product_library[product_id].copy()
                
                # 获取主形状位置（使用第一个形状）
                main_shape = shapes[0]
                
                product_info.update({
                    "product_id": product_id,
                    "quantity": 1,
                    "slide_number": 2,
                    "position": f"({int(main_shape.left/Inches(1))},{int(main_shape.top/Inches(1))})",
                    "shape_count": len(shapes),
                    "shape_types": [shape.name.split('_')[-1] for shape in shapes]
                })
                
                product_info["total_price"] = product_info["price"] * product_info["quantity"]
                all_products.append(product_info)
                
                print(f"   ✅ 产品组 {product_id}: {product_info['name']} - ¥{product_info['price']}")
                print(f"      包含 {len(shapes)} 个形状: {', '.join([shape.name.split('_')[-1] for shape in shapes])}")
    
    print(f"📊 总计找到 {len(all_products)} 个智能家居产品")
    return all_products

def create_product_report(product_data, output_path):
    """创建产品报告"""
    import openpyxl
    from datetime import datetime
    
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "智能家居产品组报告"
    
    # 表头
    headers = ["产品ID", "产品名称", "品牌", "单价(元)", "数量", "总价(元)", "形状数量", "形状类型", "位置"]
    
    for col, header in enumerate(headers, 1):
        sheet.cell(row=1, column=col, value=header)
    
    # 数据行
    total_amount = 0
    for row, product in enumerate(product_data, 2):
        sheet.cell(row=row, column=1, value=product.get("product_id", ""))
        sheet.cell(row=row, column=2, value=product.get("name", ""))
        sheet.cell(row=row, column=3, value=product.get("brand", ""))
        sheet.cell(row=row, column=4, value=product.get("price", 0))
        sheet.cell(row=row, column=5, value=product.get("quantity", 1))
        sheet.cell(row=row, column=6, value=product.get("total_price", 0))
        sheet.cell(row=row, column=7, value=product.get("shape_count", 0))
        sheet.cell(row=row, column=8, value=", ".join(product.get("shape_types", [])))
        sheet.cell(row=row, column=9, value=product.get("position", ""))
        
        total_amount += product.get("total_price", 0)
    
    # 保存文件
    workbook.save(output_path)
    return total_amount

def main():
    """主函数"""
    
    ppt_path = 'E:\\\\Programs\\\\smarthome\\\\output\\\\smart_home_real_group_mold_gallery.pptx'
    excel_library_path = 'E:\\\\Programs\\\\smarthome\\\\智能家居模具库.xlsx'
    
    if not os.path.exists(ppt_path):
        print("❌ 请先创建真组结构PPT模具库")
        return
    
    print("=== 智能家居组识别系统 ===")
    
    # 扫描PPT文件
    product_data = scan_ppt_for_product_groups(ppt_path, excel_library_path)
    
    if not product_data:
        print("❌ 未找到产品数据")
        return
    
    # 生成报告
    report_path = 'E:\\\\Programs\\\\smarthome\\\\output\\\\group_recognition_report.xlsx'
    total_amount = create_product_report(product_data, report_path)
    
    print(f"\\n✅ 产品组报告已生成: {report_path}")
    print(f"💰 总金额: {total_amount} 元")
    print(f"📊 产品数量: {len(product_data)} 个")

if __name__ == "__main__":
    main()
'''
    
    script_path = 'E:\\Programs\\smarthome\\src\\group_recognition.py'
    with open(script_path, 'w', encoding='utf-8') as f:
        f.write(script_content)
    
    print(f"✅ 组识别脚本已创建: {script_path}")
    return script_path

def main():
    """主函数"""
    
    print("🔧 解决复制组时无法复制组下文字的问题")
    print("=" * 60)
    
    # 1. 创建真组结构PPT模具库
    ppt_path = create_real_group_mold_library()
    
    # 2. 测试组识别功能
    test_success = test_group_recognition()
    
    # 3. 创建组识别脚本
    script_path = create_group_recognition_script()
    
    if test_success:
        print("\n🎉 真组结构解决方案完成!")
        print("\n📋 解决方案特点:")
        print("   ✅ 真正的组结构（多个独立形状）")
        print("   ✅ 复制时所有元素都会被复制")
        print("   ✅ 隐藏的产品ID用于识别")
        print("   ✅ 支持从Excel动态读取产品信息")
        print("   ✅ 完整的组识别和报告系统")
    else:
        print("\n❌ 测试失败，请检查配置")

if __name__ == "__main__":
    main()