#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
智能分析全屋智能方案PPT中的设备信息
"""

import os
import pandas as pd
from pptx import Presentation
import re

def smart_analyze_smart_home_plan(ppt_file_path):
    """智能分析智能家居方案PPT中的设备"""
    
    print("🔍 开始智能分析全屋智能方案...")
    print(f"📄 文件: {os.path.basename(ppt_file_path)}")
    
    # 检查文件是否存在
    if not os.path.exists(ppt_file_path):
        print(f"❌ 文件不存在: {ppt_file_path}")
        return None
    
    # 读取Excel模具库
    excel_path = 'E:\\Programs\\smarthome\\智能家居模具库.xlsx'
    if not os.path.exists(excel_path):
        print("❌ Excel模具库文件不存在")
        return None
    
    df = pd.read_excel(excel_path)
    print(f"📊 模具库包含 {len(df)} 个产品")
    
    # 创建产品映射
    product_mapping = {}
    for _, row in df.iterrows():
        product_id = row['产品ID']
        product_mapping[product_id] = {
            '设备名称': row['设备名称'],
            '品牌': row['品牌'],
            '主规格': row['主规格'],
            '设备品类': row['设备品类'],
            '单价': row['单价'],
            '设备简称': row['设备简称']
        }
    
    # 创建关键词映射（通过设备简称和关键词匹配）
    keyword_mapping = {}
    for product_id, info in product_mapping.items():
        # 使用设备简称作为关键词
        if pd.notna(info['设备简称']):
            keyword_mapping[info['设备简称'].lower()] = product_id
        
        # 使用设备名称中的关键词
        device_name = info['设备名称'].lower()
        if '开关' in device_name:
            keyword_mapping['开关'] = product_id
        if '插座' in device_name:
            keyword_mapping['插座'] = product_id
        if '传感器' in device_name:
            keyword_mapping['传感器'] = product_id
    
    # 读取PPT文件
    prs = Presentation(ppt_file_path)
    
    print(f"📊 PPT包含 {len(prs.slides)} 张幻灯片")
    
    # 统计设备信息
    device_count = {}
    total_devices = 0
    
    # 遍历所有幻灯片
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n📋 分析第 {slide_num} 张幻灯片...")
        
        # 遍历所有形状
        for shape in slide.shapes:
            # 检查是否是组
            if shape.shape_type == 6:  # GroupShape
                print(f"   🔍 发现组: {shape.name}")
                
                # 检查组内是否有设备信息
                group_texts = []
                
                # 遍历组内形状
                for sub_shape in shape.shapes:
                    if hasattr(sub_shape, 'text') and sub_shape.text:
                        text = sub_shape.text.strip()
                        if text:
                            group_texts.append(text)
                
                # 分析组内文本内容
                if group_texts:
                    print(f"      组内文本: {group_texts}")
                    
                    # 尝试通过关键词匹配设备
                    matched_product = None
                    for text in group_texts:
                        text_lower = text.lower()
                        
                        # 检查是否包含设备关键词
                        for keyword, product_id in keyword_mapping.items():
                            if keyword in text_lower:
                                matched_product = product_id
                                print(f"      ✅ 通过关键词 '{keyword}' 匹配到设备: {product_id}")
                                break
                        
                        if matched_product:
                            break
                    
                    # 如果匹配到设备，进行统计
                    if matched_product and matched_product in product_mapping:
                        product_info = product_mapping[matched_product]
                        
                        if matched_product not in device_count:
                            device_count[matched_product] = {
                                '设备名称': product_info['设备名称'],
                                '品牌': product_info['品牌'],
                                '主规格': product_info['主规格'],
                                '设备品类': product_info['设备品类'],
                                '单价': product_info['单价'],
                                '数量': 0
                            }
                        
                        device_count[matched_product]['数量'] += 1
                        total_devices += 1
                        
                        print(f"      📊 统计设备: {product_info['设备名称']} ({product_info['品牌']})")
            
            # 检查独立形状
            elif hasattr(shape, 'text') and shape.text:
                text = shape.text.strip()
                if text:
                    text_lower = text.lower()
                    
                    # 检查是否包含设备关键词
                    matched_product = None
                    for keyword, product_id in keyword_mapping.items():
                        if keyword in text_lower:
                            matched_product = product_id
                            print(f"   ✅ 通过关键词 '{keyword}' 匹配到独立设备: {product_id}")
                            break
                    
                    # 如果匹配到设备，进行统计
                    if matched_product and matched_product in product_mapping:
                        product_info = product_mapping[matched_product]
                        
                        if matched_product not in device_count:
                            device_count[matched_product] = {
                                '设备名称': product_info['设备名称'],
                                '品牌': product_info['品牌'],
                                '主规格': product_info['主规格'],
                                '设备品类': product_info['设备品类'],
                                '单价': product_info['单价'],
                                '数量': 0
                            }
                        
                        device_count[matched_product]['数量'] += 1
                        total_devices += 1
                        
                        print(f"   📊 统计独立设备: {product_info['设备名称']} ({product_info['品牌']})")
    
    return device_count, total_devices

def generate_smart_report(device_count, total_devices):
    """生成智能设备统计报告"""
    
    print("\n" + "="*60)
    print("📊 全屋智能方案设备智能分析报告")
    print("="*60)
    
    if not device_count:
        print("❌ 未识别到任何设备")
        print("\n💡 建议:")
        print("   1. 确保PPT中的设备使用了正确的产品ID格式")
        print("   2. 或者确保设备文本中包含设备关键词（如'开关'、'插座'等）")
        return
    
    # 按设备品类分组统计
    category_stats = {}
    total_cost = 0
    
    for product_id, info in device_count.items():
        category = info['设备品类']
        if category not in category_stats:
            category_stats[category] = []
        
        category_stats[category].append(info)
        total_cost += info['单价'] * info['数量']
    
    # 输出统计结果
    print(f"\n📈 总体统计:")
    print(f"   • 设备总数: {total_devices} 个")
    print(f"   • 设备种类: {len(device_count)} 种")
    print(f"   • 设备品类: {len(category_stats)} 类")
    print(f"   • 预估总价: {total_cost:.2f} 元")
    
    # 按品类输出详细信息
    for category, devices in category_stats.items():
        print(f"\n🏷️  {category}:")
        
        for device in devices:
            print(f"   📋 {device['设备名称']}")
            print(f"      • 品牌: {device['品牌']}")
            print(f"      • 规格: {device['主规格']}")
            print(f"      • 数量: {device['数量']} 个")
            print(f"      • 单价: {device['单价']} 元")
            print(f"      • 小计: {device['单价'] * device['数量']:.2f} 元")
    
    print(f"\n💰 总金额: {total_cost:.2f} 元")
    
    # 生成Excel报告
    report_data = []
    for product_id, info in device_count.items():
        report_data.append({
            '产品ID': product_id,
            '设备品类': info['设备品类'],
            '设备名称': info['设备名称'],
            '品牌': info['品牌'],
            '主规格': info['主规格'],
            '单价': info['单价'],
            '数量': info['数量'],
            '小计': info['单价'] * info['数量']
        })
    
    report_df = pd.DataFrame(report_data)
    report_path = 'E:\\Programs\\smarthome\\output\\全屋智能方案智能分析报告.xlsx'
    report_df.to_excel(report_path, index=False)
    
    print(f"\n📄 详细报告已保存到: {os.path.basename(report_path)}")

def main():
    """主函数"""
    
    ppt_file_path = 'E:\\Programs\\smarthome\\全屋智能方案.pptx'
    
    print("🔧 全屋智能方案智能分析工具")
    print("="*60)
    
    # 智能分析PPT文件
    device_count, total_devices = smart_analyze_smart_home_plan(ppt_file_path)
    
    if device_count is not None:
        # 生成报告
        generate_smart_report(device_count, total_devices)
    else:
        print("❌ 分析失败，请检查文件路径和格式")

if __name__ == "__main__":
    main()