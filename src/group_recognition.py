#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
智能家居组识别脚本
支持识别真组结构中的产品信息
"""

from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import os

def read_excel_product_library(excel_path):
    """从Excel模具库读取产品信息"""
    if not os.path.exists(excel_path):
        print(f"❌ Excel模具库文件不存在: {excel_path}")
        return {}
    
    df = pd.read_excel(excel_path)
    product_library = {}
    
    for index, row in df.iterrows():
        product_id = row.get('产品ID', '')
        if pd.notna(product_id) and product_id:
            product_library[product_id] = {
                "name": row.get('设备名称', ''),
                "price": int(row.get('单价', 0)) if pd.notna(row.get('单价')) else 0,
                "brand": row.get('品牌', ''),
                "model": row.get('主规格', ''),
                "category": row.get('设备品类', '')
            }
    
    print(f"✅ 从Excel读取了 {len(product_library)} 个产品信息")
    return product_library

def extract_product_id_from_shape_name(shape_name):
    """从形状名称提取产品ID"""
    if not shape_name or 'smart_home_' not in shape_name:
        return None
    
    # 格式: smart_home_switch_1_lp_id
    parts = shape_name.split('smart_home_')[1].split('_')
    if len(parts) >= 3:
        return f"{parts[0]}_{parts[1]}_{parts[2]}"
    elif len(parts) >= 2:
        return f"{parts[0]}_{parts[1]}"
    
    return None

def scan_ppt_for_product_groups(ppt_path, excel_library_path):
    """扫描PPT文件中的产品组"""
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")
    
    # 读取产品库
    product_library = read_excel_product_library(excel_library_path)
    
    prs = Presentation(ppt_path)
    all_products = []
    
    print(f"🔍 扫描PPT文件: {ppt_path}")
    print(f"📊 幻灯片数量: {len(prs.slides)}")
    
    # 处理第二张幻灯片（产品模具页）
    if len(prs.slides) > 1:
        slide = prs.slides[1]
        
        # 按产品ID分组形状
        product_groups = {}
        for shape in slide.shapes:
            if hasattr(shape, 'name') and shape.name:
                product_id = extract_product_id_from_shape_name(shape.name)
                if product_id:
                    if product_id not in product_groups:
                        product_groups[product_id] = []
                    product_groups[product_id].append(shape)
        
        print(f"📦 识别到 {len(product_groups)} 个产品组")
        
        # 处理每个产品组
        for product_id, shapes in product_groups.items():
            if product_id in product_library:
                product_info = product_library[product_id].copy()
                
                # 获取主形状位置（使用第一个形状）
                main_shape = shapes[0]
                
                product_info.update({
                    "product_id": product_id,
                    "quantity": 1,
                    "slide_number": 2,
                    "position": f"({int(main_shape.left/Inches(1))},{int(main_shape.top/Inches(1))})",
                    "shape_count": len(shapes),
                    "shape_types": [shape.name.split('_')[-1] for shape in shapes]
                })
                
                product_info["total_price"] = product_info["price"] * product_info["quantity"]
                all_products.append(product_info)
                
                print(f"   ✅ 产品组 {product_id}: {product_info['name']} - ¥{product_info['price']}")
                print(f"      包含 {len(shapes)} 个形状: {', '.join([shape.name.split('_')[-1] for shape in shapes])}")
    
    print(f"📊 总计找到 {len(all_products)} 个智能家居产品")
    return all_products

def create_product_report(product_data, output_path):
    """创建产品报告"""
    import openpyxl
    from datetime import datetime
    
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "智能家居产品组报告"
    
    # 表头
    headers = ["产品ID", "产品名称", "品牌", "单价(元)", "数量", "总价(元)", "形状数量", "形状类型", "位置"]
    
    for col, header in enumerate(headers, 1):
        sheet.cell(row=1, column=col, value=header)
    
    # 数据行
    total_amount = 0
    for row, product in enumerate(product_data, 2):
        sheet.cell(row=row, column=1, value=product.get("product_id", ""))
        sheet.cell(row=row, column=2, value=product.get("name", ""))
        sheet.cell(row=row, column=3, value=product.get("brand", ""))
        sheet.cell(row=row, column=4, value=product.get("price", 0))
        sheet.cell(row=row, column=5, value=product.get("quantity", 1))
        sheet.cell(row=row, column=6, value=product.get("total_price", 0))
        sheet.cell(row=row, column=7, value=product.get("shape_count", 0))
        sheet.cell(row=row, column=8, value=", ".join(product.get("shape_types", [])))
        sheet.cell(row=row, column=9, value=product.get("position", ""))
        
        total_amount += product.get("total_price", 0)
    
    # 保存文件
    workbook.save(output_path)
    return total_amount

def main():
    """主函数"""
    
    ppt_path = 'E:\\Programs\\smarthome\\output\\smart_home_real_group_mold_gallery.pptx'
    excel_library_path = 'E:\\Programs\\smarthome\\智能家居模具库.xlsx'
    
    if not os.path.exists(ppt_path):
        print("❌ 请先创建真组结构PPT模具库")
        return
    
    print("=== 智能家居组识别系统 ===")
    
    # 扫描PPT文件
    product_data = scan_ppt_for_product_groups(ppt_path, excel_library_path)
    
    if not product_data:
        print("❌ 未找到产品数据")
        return
    
    # 生成报告
    report_path = 'E:\\Programs\\smarthome\\output\\group_recognition_report.xlsx'
    total_amount = create_product_report(product_data, report_path)
    
    print(f"\n✅ 产品组报告已生成: {report_path}")
    print(f"💰 总金额: {total_amount} 元")
    print(f"📊 产品数量: {len(product_data)} 个")

if __name__ == "__main__":
    main()
