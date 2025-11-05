#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT模具库改进模块 - 任务3：创建PPT模具库改进模块
为PPT中的设备组添加pdid标签
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import json
import os
from typing import Dict, List, Tuple

class PPTEnhancer:
    """PPT模具库改进器"""
    
    def __init__(self, ppt_path: str, mapping_path: str = None):
        """
        初始化PPT改进器
        
        Args:
            ppt_path: PPT文件路径
            mapping_path: 产品ID映射表路径
        """
        self.ppt_path = ppt_path
        self.mapping_path = mapping_path
        self.presentation = None
        self.mapping = {}
        
    def load_presentation(self) -> bool:
        """
        加载PPT文件
        
        Returns:
            bool: 是否成功加载
        """
        try:
            self.presentation = Presentation(self.ppt_path)
            print(f"✅ 成功加载PPT文件: {self.ppt_path}")
            print(f"📊 幻灯片数量: {len(self.presentation.slides)}")
            return True
        except Exception as e:
            print(f"❌ 加载PPT文件失败: {e}")
            return False
    
    def load_mapping_table(self) -> bool:
        """
        加载产品ID映射表
        
        Returns:
            bool: 是否成功加载
        """
        if not self.mapping_path or not os.path.exists(self.mapping_path):
            print("⚠️ 未提供映射表路径或文件不存在，将使用默认映射")
            return True
        
        try:
            with open(self.mapping_path, 'r', encoding='utf-8') as f:
                self.mapping = json.load(f)
            print(f"✅ 成功加载产品ID映射表: {self.mapping_path}")
            print(f"📋 映射关系: {self.mapping}")
            return True
        except Exception as e:
            print(f"❌ 加载映射表失败: {e}")
            return False
    
    def analyze_slides(self) -> Dict[int, List[Dict]]:
        """
        分析幻灯片中的设备组
        
        Returns:
            Dict[int, List[Dict]]: 幻灯片索引到设备组信息的映射
        """
        if self.presentation is None:
            return {}
        
        device_groups = {}
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            print(f"\n📋 分析第{slide_idx + 1}张幻灯片:")
            
            # 查找设备组
            device_shapes = []
            for shape in slide.shapes:
                shape_info = {
                    'shape': shape,
                    'name': shape.name,
                    'type': type(shape).__name__,
                    'has_text': shape.has_text_frame,
                    'text': shape.text if shape.has_text_frame else ""
                }
                
                # 判断是否为设备组相关形状
                if 'smart_home_switch' in shape.name.lower() or 'switch' in shape.name.lower():
                    device_shapes.append(shape_info)
                    print(f"   🔍 发现设备组形状: {shape.name}")
                elif shape.has_text_frame and ('开关' in shape.text or 'switch' in shape.text.lower()):
                    device_shapes.append(shape_info)
                    print(f"   🔍 发现设备组文本: {shape.text[:30]}...")
            
            device_groups[slide_idx] = device_shapes
            print(f"   📊 本页设备组数量: {len(device_shapes)}")
        
        return device_groups
    
    def get_device_pdid(self, shape_info: Dict) -> int:
        """
        根据形状信息获取对应的产品ID
        
        Args:
            shape_info: 形状信息字典
            
        Returns:
            int: 产品ID，如果无法确定返回0
        """
        shape_name = shape_info['name'].lower()
        shape_text = shape_info['text'].lower()
        
        # 根据形状名称匹配产品ID
        if 'switch_1' in shape_name:
            return 1 if 'lp' in shape_name else 5
        elif 'switch_2' in shape_name:
            return 2 if 'lp' in shape_name else 6
        elif 'switch_3' in shape_name:
            return 3 if 'lp' in shape_name else 7
        elif 'switch_4' in shape_name:
            return 4 if 'lp' in shape_name else 8
        
        # 根据文本内容匹配
        if '一键' in shape_text or '1键' in shape_text:
            return 1 if '领普' in shape_text else 5
        elif '二键' in shape_text or '2键' in shape_text:
            return 2 if '领普' in shape_text else 6
        elif '三键' in shape_text or '3键' in shape_text:
            return 3 if '领普' in shape_text else 7
        elif '四键' in shape_text or '4键' in shape_text:
            return 4 if '领普' in shape_text else 8
        
        return 0
    
    def add_pdid_label(self, shape_info: Dict, pdid: int) -> bool:
        """
        为设备组添加pdid标签
        
        Args:
            shape_info: 形状信息字典
            pdid: 产品ID
            
        Returns:
            bool: 是否成功添加
        """
        shape = shape_info['shape']
        
        try:
            # 获取幻灯片对象 - 需要遍历幻灯片来找到包含该形状的幻灯片
            slide = None
            for slide_idx, current_slide in enumerate(self.presentation.slides):
                for slide_shape in current_slide.shapes:
                    if slide_shape == shape:
                        slide = current_slide
                        break
                if slide:
                    break
            
            if slide is None:
                print(f"   ❌ 无法找到包含形状 {shape_info['name']} 的幻灯片")
                return False
            
            # 获取设备组的位置和尺寸
            left = shape.left
            top = shape.top + shape.height
            width = shape.width
            height = Inches(0.3)  # 标签高度
            
            # 创建pdid标签文本框
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.name = f"pdid_label_{pdid}"
            
            # 设置文本框样式
            text_frame = textbox.text_frame
            text_frame.clear()  # 清除默认文本
            
            # 添加pdid文本
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = f"pdid: {pdid}"
            
            # 设置文本格式
            font = run.font
            font.name = 'Arial'
            font.size = Pt(6)  # 最小号字体
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)  # 黑色文字
            
            # 设置文本框背景 - 透明不填充
            fill = textbox.fill
            fill.background()  # 透明背景
            
            # 设置文本框边框 - 无边框
            line = textbox.line
            line.fill.background()  # 透明边框
            
            # 设置文本对齐
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p.alignment = 1  # 居中对齐
            
            print(f"   ✅ 为设备组 {shape_info['name']} 添加pdid标签: {pdid}")
            return True
            
        except Exception as e:
            print(f"   ❌ 添加pdid标签失败: {e}")
            return False
    
    def enhance_presentation(self) -> bool:
        """
        改进PPT模具库
        
        Returns:
            bool: 是否成功改进
        """
        if self.presentation is None:
            print("❌ 请先加载PPT文件")
            return False
        
        # 分析幻灯片
        device_groups = self.analyze_slides()
        
        if not device_groups:
            print("❌ 未发现设备组")
            return False
        
        total_added = 0
        
        # 为每个设备组添加pdid标签
        for slide_idx, groups in device_groups.items():
            print(f"\n🎯 处理第{slide_idx + 1}张幻灯片:")
            
            for shape_info in groups:
                pdid = self.get_device_pdid(shape_info)
                if pdid > 0:
                    if self.add_pdid_label(shape_info, pdid):
                        total_added += 1
                else:
                    print(f"   ⚠️ 无法确定设备组 {shape_info['name']} 的产品ID")
        
        print(f"\n✅ 共添加了 {total_added} 个pdid标签")
        return total_added > 0
    
    def save_enhanced_ppt(self, output_path: str = None) -> bool:
        """
        保存改进后的PPT文件
        
        Args:
            output_path: 输出文件路径，如果为None则覆盖原文件
            
        Returns:
            bool: 是否成功保存
        """
        if self.presentation is None:
            print("❌ 没有PPT数据可保存")
            return False
        
        if output_path is None:
            output_path = self.ppt_path
        
        try:
            self.presentation.save(output_path)
            print(f"✅ 改进后的PPT文件已保存: {output_path}")
            return True
        except Exception as e:
            print(f"❌ 保存PPT文件失败: {e}")
            return False

def enhance_ppt_library(ppt_path: str, mapping_path: str = None, output_path: str = None) -> bool:
    """
    改进PPT模具库的主函数
    
    Args:
        ppt_path: 输入PPT文件路径
        mapping_path: 产品ID映射表路径
        output_path: 输出PPT文件路径
        
    Returns:
        bool: 是否成功改进
    """
    print("=" * 60)
    print("🔧 开始PPT模具库改进 - 任务3")
    print("=" * 60)
    
    # 初始化PPT改进器
    enhancer = PPTEnhancer(ppt_path, mapping_path)
    
    # 加载PPT文件
    if not enhancer.load_presentation():
        return False
    
    # 加载映射表（可选）
    if not enhancer.load_mapping_table():
        return False
    
    # 改进PPT
    if not enhancer.enhance_presentation():
        return False
    
    # 保存结果
    if not enhancer.save_enhanced_ppt(output_path):
        return False
    
    print("=" * 60)
    print("✅ PPT模具库改进任务完成")
    print("=" * 60)
    
    return True

if __name__ == "__main__":
    # 测试函数
    ppt_path = "E:\\Programs\\smarthome\\智能家居模具库.pptx"
    success = enhance_ppt_library(ppt_path)
    
    if success:
        print("🎯 PPT模具库改进成功")
    else:
        print("❌ PPT模具库改进失败")