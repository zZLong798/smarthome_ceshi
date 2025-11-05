#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
增强的智能家居组识别脚本
支持识别真组结构中的产品信息（组名称=产品ID）
"""

from pptx import Presentation
import pandas as pd
import os

def read_excel_product_library(excel_path):
    """从Excel模具库读取产品信息"""
    if not os.path.exists(excel_path):
        print(f"❌ Excel模具库文件不存在: {excel_path}")
        return {}
    
    df = pd.read_excel(excel_path)
    product_library = {}
    
    for index, row in df.iterrows():
        product_id = row.get('产品ID', '')
        if pd.notna(product_id) and product_id:
            product_library[product_id] = {
                '设备名称': row.get('设备名称', ''),
                '品牌': row.get('品牌', ''),
                '主规格': row.get('主规格', ''),
                '设备品类': row.get('设备品类', ''),
                '单价': row.get('单价', 0),
                '设备简称': row.get('设备简称', '')
            }
    
    print(f"📊 从Excel读取 {len(product_library)} 个产品信息")
    return product_library

def analyze_true_group_ppt(ppt_path, product_library):
    """分析真组结构PPT文件"""
    
    if not os.path.exists(ppt_path):
        print(f"❌ PPT文件不存在: {ppt_path}")
        return {}
    
    prs = Presentation(ppt_path)
    print(f"📄 分析PPT文件: {os.path.basename(ppt_path)}")
    print(f"📊 幻灯片数量: {len(prs.slides)}")
    
    device_count = {}
    total_devices = 0
    
    # 遍历所有幻灯片
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n📋 分析第 {slide_num} 张幻灯片...")
        
        # 统计产品组
        product_groups = {}
        for shape in slide.shapes:
            if hasattr(shape, 'name') and shape.name:
                # 通过_id后缀识别产品组
                if shape.name.endswith('_id'):
                    product_id = shape.name.replace('_id', '')
                    if product_id in product_library:
                        if product_id not in product_groups:
                            product_groups[product_id] = 0
                        product_groups[product_id] += 1
        
        # 统计设备数量
        for product_id, count in product_groups.items():
            if product_id not in device_count:
                device_count[product_id] = 0
            device_count[product_id] += count
            total_devices += count
            
            product_info = product_library[product_id]
            print(f"   ✅ 识别到设备: {product_info['设备名称']} ({product_info['品牌']}) x{count}")
    
    return device_count, total_devices

def generate_true_group_report(device_count, total_devices, product_library):
    """生成真组结构报告"""
    
    print("\n" + "="*60)
    print("📊 真组结构设备统计报告")
    print("="*60)
    
    if not device_count:
        print("❌ 未识别到任何设备")
        return
    
    # 按设备品类分组统计
    category_stats = {}
    total_cost = 0
    
    for product_id, count in device_count.items():
        if product_id in product_library:
            product_info = product_library[product_id]
            category = product_info['设备品类']
            
            if category not in category_stats:
                category_stats[category] = []
            
            category_stats[category].append({
                '设备名称': product_info['设备名称'],
                '品牌': product_info['品牌'],
                '主规格': product_info['主规格'],
                '单价': product_info['单价'],
                '数量': count
            })
            
            total_cost += product_info['单价'] * count
    
    # 输出统计结果
    print(f"\n📈 总体统计:")
    print(f"   • 设备总数: {total_devices} 个")
    print(f"   • 设备种类: {len(device_count)} 种")
    print(f"   • 设备品类: {len(category_stats)} 类")
    print(f"   • 预估总价: {total_cost:.2f} 元")
    
    # 按品类输出详细信息
    for category, devices in category_stats.items():
        print(f"\n🏷️  {category}:")
        
        for device in devices:
            print(f"   📋 {device['设备名称']}")
            print(f"      • 品牌: {device['品牌']}")
            print(f"      • 规格: {device['主规格']}")
            print(f"      • 数量: {device['数量']} 个")
            print(f"      • 单价: {device['单价']} 元")
            print(f"      • 小计: {device['单价'] * device['数量']:.2f} 元")
    
    print(f"\n💰 总金额: {total_cost:.2f} 元")
    
    # 生成Excel报告
    report_data = []
    for product_id, count in device_count.items():
        if product_id in product_library:
            product_info = product_library[product_id]
            report_data.append({
                '产品ID': product_id,
                '设备品类': product_info['设备品类'],
                '设备名称': product_info['设备名称'],
                '品牌': product_info['品牌'],
                '主规格': product_info['主规格'],
                '单价': product_info['单价'],
                '数量': count,
                '小计': product_info['单价'] * count
            })
    
    report_df = pd.DataFrame(report_data)
    report_path = 'E:\\Programs\\smarthome\\output\\true_group_recognition_report.xlsx'
    report_df.to_excel(report_path, index=False)
    
    print(f"\n📄 详细报告已保存到: {os.path.basename(report_path)}")

def main():
    """主函数"""
    
    # 读取Excel模具库
    excel_path = 'E:\\Programs\\smarthome\\智能家居模具库.xlsx'
    product_library = read_excel_product_library(excel_path)
    
    if not product_library:
        print("❌ 无法读取产品库信息")
        return
    
    # 分析全屋智能方案PPT
    ppt_path = 'E:\\Programs\\smarthome\\全屋智能方案.pptx'
    device_count, total_devices = analyze_true_group_ppt(ppt_path, product_library)
    
    # 生成报告
    generate_true_group_report(device_count, total_devices, product_library)

if __name__ == "__main__":
    main()
