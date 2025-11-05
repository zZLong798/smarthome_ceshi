#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
分析原始PPT文件中的设备信息
通过产品ID识别设备名称、主规格和数量
"""

import os
import pandas as pd
from pptx import Presentation

def load_product_library():
    """加载智能家居模具库"""
    
    excel_file_path = 'E:\\Programs\\smarthome\\智能家居模具库.xlsx'
    
    if not os.path.exists(excel_file_path):
        print(f"❌ 模具库文件不存在: {excel_file_path}")
        return None
    
    try:
        df = pd.read_excel(excel_file_path)
        print(f"✅ 成功加载模具库，共 {len(df)} 个产品")
        return df
    except Exception as e:
        print(f"❌ 加载模具库失败: {e}")
        return None

def analyze_original_ppt_devices(ppt_file_path, product_df):
    """分析原始PPT中的设备"""
    
    print("\n🔍 分析原始PPT中的设备信息")
    print("="*60)
    
    if not os.path.exists(ppt_file_path):
        print(f"❌ PPT文件不存在: {ppt_file_path}")
        return {}
    
    prs = Presentation(ppt_file_path)
    
    # 收集所有组信息
    all_groups = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, 'name') and shape.name and hasattr(shape, 'shapes') and shape.shapes:
                all_groups.append({
                    'slide': slide_num,
                    'name': shape.name,
                    'shape': shape
                })
    
    print(f"📊 发现 {len(all_groups)} 个组")
    
    # 统计设备数量
    device_counts = {}
    
    # 方法1：通过产品ID识别
    product_id_groups = []
    for group in all_groups:
        if '_id' in group['name']:
            product_id = group['name'].replace('_id', '')
            product_id_groups.append({
                'slide': group['slide'],
                'product_id': product_id,
                'group_name': group['name']
            })
    
    print(f"🎯 通过产品ID识别到 {len(product_id_groups)} 个设备组")
    
    # 方法2：通过组名称关键词识别
    keyword_groups = []
    for group in all_groups:
        group_name = group['name'].lower()
        if 'switch' in group_name or 'smart' in group_name:
            keyword_groups.append(group)
    
    print(f"🔍 通过关键词识别到 {len(keyword_groups)} 个疑似设备组")
    
    # 优先使用产品ID识别
    for group in product_id_groups:
        product_id = group['product_id']
        
        # 在模具库中查找产品信息
        product_info = product_df[product_df['产品ID'] == product_id]
        
        if not product_info.empty:
            device_name = product_info.iloc[0]['设备名称']
            brand = product_info.iloc[0]['品牌']
            specification = product_info.iloc[0]['主规格']
            
            # 创建设备标识
            device_key = f"{brand}_{device_name}_{specification}"
            
            if device_key not in device_counts:
                device_counts[device_key] = {
                    'brand': brand,
                    'device_name': device_name,
                    'specification': specification,
                    'product_id': product_id,
                    'count': 0,
                    'recognition_method': '产品ID'
                }
            
            device_counts[device_key]['count'] += 1
        else:
            print(f"⚠️ 未找到产品ID '{product_id}' 对应的设备信息")
    
    # 如果没有通过产品ID识别到设备，尝试关键词识别
    if not device_counts and keyword_groups:
        print("🔧 尝试通过组名称关键词识别设备...")
        
        # 创建组名称到产品ID的映射
        name_mapping = {
            'smart_home_switch_1': 'switch_1_yl',
            'smart_home_switch_2': 'switch_2_yl',
            'smart_home_switch_3': 'switch_3_yl',
            'smart_home_switch_4': 'switch_4_yl'
        }
        
        for group in keyword_groups:
            group_name = group['name']
            
            # 检查是否有映射关系
            if group_name in name_mapping:
                product_id = name_mapping[group_name]
                
                # 在模具库中查找产品信息
                product_info = product_df[product_df['产品ID'] == product_id]
                
                if not product_info.empty:
                    device_name = product_info.iloc[0]['设备名称']
                    brand = product_info.iloc[0]['品牌']
                    specification = product_info.iloc[0]['主规格']
                    
                    # 创建设备标识
                    device_key = f"{brand}_{device_name}_{specification}"
                    
                    if device_key not in device_counts:
                        device_counts[device_key] = {
                            'brand': brand,
                            'device_name': device_name,
                            'specification': specification,
                            'product_id': product_id,
                            'count': 0,
                            'recognition_method': '组名称映射'
                        }
                    
                    device_counts[device_key]['count'] += 1
                    print(f"   📍 映射 '{group_name}' -> '{product_id}'")
    
    return device_counts

def generate_device_report(device_counts):
    """生成设备报告"""
    
    print("\n📊 设备统计报告")
    print("="*60)
    
    if not device_counts:
        print("❌ 未发现任何设备")
        return
    
    total_devices = sum(device['count'] for device in device_counts.values())
    print(f"📈 总共发现 {total_devices} 个设备")
    print()
    
    # 按品牌和设备类型排序显示
    sorted_devices = sorted(device_counts.items(), key=lambda x: (x[1]['brand'], x[1]['device_name']))
    
    for device_key, device_info in sorted_devices:
        print(f"🏷️ 品牌: {device_info['brand']}")
        print(f"📱 设备: {device_info['device_name']}")
        print(f"🔧 主规格: {device_info['specification']}")
        print(f"🆔 产品ID: {device_info['product_id']}")
        print(f"📦 数量: {device_info['count']} 个")
        print(f"🔍 识别方式: {device_info['recognition_method']}")
        print("-" * 40)

def main():
    """主函数"""
    
    # 使用原始PPT文件
    ppt_file_path = 'E:\\Programs\\smarthome\\全屋智能方案.pptx'
    
    print("🔧 原始PPT设备分析工具")
    print("="*60)
    
    # 加载模具库
    product_df = load_product_library()
    if product_df is None:
        return
    
    # 分析PPT设备
    device_counts = analyze_original_ppt_devices(ppt_file_path, product_df)
    
    # 生成报告
    generate_device_report(device_counts)
    
    print("\n" + "="*60)
    print("✅ 分析完成")

if __name__ == "__main__":
    main()