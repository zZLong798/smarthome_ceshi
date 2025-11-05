#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
为PPT中的设备添加产品ID
帮助用户正确设置产品ID格式
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

def show_product_id_guide():
    """显示产品ID使用指南"""
    
    print("📋 产品ID使用指南")
    print("="*60)
    print("🔧 正确设置产品ID的方法:")
    print("")
    print("1. 在PPT中为每个设备组添加一个隐藏的形状")
    print("2. 形状名称使用产品ID + '_id'后缀")
    print("3. 例如: switch_1_lp_id, switch_4_yl_id")
    print("")
    print("📋 模具库中的产品ID列表:")
    
    # 读取Excel模具库
    excel_path = 'E:\\Programs\\smarthome\\智能家居模具库.xlsx'
    df = pd.read_excel(excel_path)
    
    for _, row in df.iterrows():
        product_id = row['产品ID']
        device_name = row['设备名称']
        brand = row['品牌']
        
        print(f"   • {product_id}_id -> {device_name} ({brand})")
    
    print("")
    print("💡 操作步骤:")
    print("   1. 在PPT中插入一个文本框或矩形")
    print("   2. 设置形状名称为产品ID + '_id'")
    print("   3. 将这个形状放在设备组内")
    print("   4. 可以设置形状为透明或隐藏")

def check_current_ppt_structure(ppt_file_path):
    """检查当前PPT的结构"""
    
    print("\n🔍 检查当前PPT结构")
    print("="*60)
    
    if not os.path.exists(ppt_file_path):
        print(f"❌ PPT文件不存在: {ppt_file_path}")
        return
    
    prs = Presentation(ppt_file_path)
    
    print(f"📊 PPT包含 {len(prs.slides)} 张幻灯片")
    
    # 统计组和形状信息
    total_groups = 0
    total_shapes = 0
    shape_names = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_groups = 0
        slide_shapes = 0
        
        for shape in slide.shapes:
            total_shapes += 1
            slide_shapes += 1
            
            if hasattr(shape, 'name') and shape.name:
                shape_names.append(shape.name)
            
            # 检查是否是组
            if hasattr(shape, 'shapes') and shape.shapes:
                total_groups += 1
                slide_groups += 1
        
        if slide_groups > 0:
            print(f"   第 {slide_num} 张幻灯片: {slide_groups} 个组, {slide_shapes} 个形状")
    
    print(f"\n📊 总体统计:")
    print(f"   • 总组数: {total_groups}")
    print(f"   • 总形状数: {total_shapes}")
    
    # 显示独特的形状名称
    unique_names = list(set(shape_names))
    print(f"   • 独特形状名称: {len(unique_names)} 个")
    
    if unique_names:
        print("\n📋 当前PPT中的形状名称:")
        for name in sorted(unique_names)[:20]:  # 只显示前20个
            print(f"   • '{name}'")

def create_sample_ppt_with_ids():
    """创建包含产品ID的示例PPT"""
    
    print("\n🎯 创建示例PPT")
    print("="*60)
    
    # 创建新的PPT
    prs = Presentation()
    
    # 添加标题幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "智能家居设备示例"
    
    # 添加内容幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "设备配置示例"
    
    # 添加示例设备组
    left = Inches(1)
    top = Inches(2)
    width = Inches(3)
    height = Inches(1)
    
    # 创建设备组
    from pptx.enum.shapes import MSO_SHAPE
    
    # 示例1: 四键易来开关
    group_shapes = slide.shapes
    
    # 主设备形状
    device_shape = group_shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    device_shape.text = "四键智能开关"
    device_shape.name = "设备形状"
    
    # 产品ID标识形状（隐藏）
    id_shape = group_shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + width - 0.2, top, 0.1, 0.1
    )
    id_shape.name = "switch_4_yl_id"  # 产品ID + _id后缀
    id_shape.fill.background()  # 设置为透明
    id_shape.line.fill.background()  # 边框透明
    
    print("✅ 创建了示例PPT，包含产品ID标识")
    print("   • 设备形状: '设备形状'")
    print("   • 产品ID形状: 'switch_4_yl_id'")
    
    # 保存示例PPT
    sample_path = 'E:\\Programs\\smarthome\\output\\产品ID示例.pptx'
    prs.save(sample_path)
    
    print(f"📄 示例PPT已保存到: {os.path.basename(sample_path)}")

def main():
    """主函数"""
    
    ppt_file_path = 'E:\\Programs\\smarthome\\全屋智能方案.pptx'
    
    print("🔧 产品ID配置工具")
    print("="*60)
    
    # 显示使用指南
    show_product_id_guide()
    
    # 检查当前PPT结构
    check_current_ppt_structure(ppt_file_path)
    
    # 创建示例PPT
    create_sample_ppt_with_ids()
    
    print("\n" + "="*60)
    print("📋 下一步操作:")
    print("="*60)
    print("1. 打开您的PPT文件")
    print("2. 为每个设备组添加一个隐藏的形状")
    print("3. 设置形状名称为产品ID + '_id'后缀")
    print("4. 保存PPT后重新运行分析工具")
    print("")
    print("💡 示例:")
    print("   • 四键领普开关: switch_4_lp_id")
    print("   • 一键易来开关: switch_1_yl_id")
    print("   • 三键领普开关: switch_3_lp_id")

if __name__ == "__main__":
    main()