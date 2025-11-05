#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
通过产品ID识别PPT中的设备
正确实现：通过产品ID后缀识别设备信息
"""

import os
import pandas as pd
from pptx import Presentation
import re

def analyze_by_product_id(ppt_file_path):
    """通过产品ID识别PPT中的设备"""
    
    print("🔍 开始通过产品ID识别设备...")
    print(f"📄 文件: {os.path.basename(ppt_file_path)}")
    
    # 检查文件是否存在
    if not os.path.exists(ppt_file_path):
        print(f"❌ 文件不存在: {ppt_file_path}")
        return None
    
    # 读取Excel模具库
    excel_path = 'E:\\Programs\\smarthome\\智能家居模具库.xlsx'
    if not os.path.exists(excel_path):
        print("❌ Excel模具库文件不存在")
        return None
    
    df = pd.read_excel(excel_path)
    print(f"📊 模具库包含 {len(df)} 个产品")
    
    # 创建产品ID映射
    product_mapping = {}
    for _, row in df.iterrows():
        product_id = row['产品ID']
        product_mapping[product_id] = {
            '设备名称': row['设备名称'],
            '品牌': row['品牌'],
            '主规格': row['主规格'],
            '设备品类': row['设备品类'],
            '单价': row['单价']
        }
    
    # 读取PPT文件
    prs = Presentation(ppt_file_path)
    
    print(f"📊 PPT包含 {len(prs.slides)} 张幻灯片")
    
    # 统计设备信息
    device_count = {}
    total_devices = 0
    
    # 遍历所有幻灯片
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n📋 分析第 {slide_num} 张幻灯片...")
        
        # 遍历所有形状
        for shape_num, shape in enumerate(slide.shapes, 1):
            
            # 检查形状是否有名称
            if hasattr(shape, 'name') and shape.name:
                shape_name = shape.name
                
                # 检查是否是组
                if hasattr(shape, 'shapes') and shape.shapes:
                    print(f"   🔍 发现组 #{shape_num}: '{shape_name}'")
                    
                    # 遍历组内形状，寻找产品ID
                    for sub_shape_num, sub_shape in enumerate(shape.shapes, 1):
                        if hasattr(sub_shape, 'name') and sub_shape.name:
                            sub_shape_name = sub_shape.name
                            
                            # 检查是否包含产品ID后缀
                            if '_id' in sub_shape_name:
                                # 提取产品ID（去掉_id后缀）
                                product_id = sub_shape_name.replace('_id', '')
                                
                                print(f"      📍 发现产品ID形状: '{sub_shape_name}'")
                                print(f"          提取的产品ID: {product_id}")
                                
                                # 检查产品ID是否在模具库中
                                if product_id in product_mapping:
                                    product_info = product_mapping[product_id]
                                    
                                    print(f"      ✅ 匹配到设备: {product_info['设备名称']} ({product_info['品牌']})")
                                    
                                    # 统计设备
                                    if product_id not in device_count:
                                        device_count[product_id] = {
                                            '设备名称': product_info['设备名称'],
                                            '品牌': product_info['品牌'],
                                            '主规格': product_info['主规格'],
                                            '设备品类': product_info['设备品类'],
                                            '单价': product_info['单价'],
                                            '数量': 0
                                        }
                                    
                                    device_count[product_id]['数量'] += 1
                                    total_devices += 1
                                    
                                    print(f"      📊 统计设备: {product_info['设备名称']} ({product_info['品牌']})")
                                else:
                                    print(f"      ❌ 产品ID '{product_id}' 不在模具库中")
                            
                            # 检查形状名称是否直接是产品ID
                            elif sub_shape_name in product_mapping:
                                product_id = sub_shape_name
                                product_info = product_mapping[product_id]
                                
                                print(f"      📍 发现产品ID形状: '{sub_shape_name}'")
                                print(f"      ✅ 匹配到设备: {product_info['设备名称']} ({product_info['品牌']})")
                                
                                # 统计设备
                                if product_id not in device_count:
                                    device_count[product_id] = {
                                        '设备名称': product_info['设备名称'],
                                        '品牌': product_info['品牌'],
                                        '主规格': product_info['主规格'],
                                        '设备品类': product_info['设备品类'],
                                        '单价': product_info['单价'],
                                        '数量': 0
                                    }
                                
                                device_count[product_id]['数量'] += 1
                                total_devices += 1
                                
                                print(f"      📊 统计设备: {product_info['设备名称']} ({product_info['品牌']})")
                
                # 检查独立形状是否包含产品ID
                elif '_id' in shape_name:
                    # 提取产品ID（去掉_id后缀）
                    product_id = shape_name.replace('_id', '')
                    
                    print(f"   📍 发现独立产品ID形状: '{shape_name}'")
                    print(f"       提取的产品ID: {product_id}")
                    
                    # 检查产品ID是否在模具库中
                    if product_id in product_mapping:
                        product_info = product_mapping[product_id]
                        
                        print(f"   ✅ 匹配到设备: {product_info['设备名称']} ({product_info['品牌']})")
                        
                        # 统计设备
                        if product_id not in device_count:
                            device_count[product_id] = {
                                '设备名称': product_info['设备名称'],
                                '品牌': product_info['品牌'],
                                '主规格': product_info['主规格'],
                                '设备品类': product_info['设备品类'],
                                '单价': product_info['单价'],
                                '数量': 0
                            }
                        
                        device_count[product_id]['数量'] += 1
                        total_devices += 1
                        
                        print(f"   📊 统计设备: {product_info['设备名称']} ({product_info['品牌']})")
                    else:
                        print(f"   ❌ 产品ID '{product_id}' 不在模具库中")
                
                # 检查形状名称是否直接是产品ID
                elif shape_name in product_mapping:
                    product_id = shape_name
                    product_info = product_mapping[product_id]
                    
                    print(f"   📍 发现独立产品ID形状: '{shape_name}'")
                    print(f"   ✅ 匹配到设备: {product_info['设备名称']} ({product_info['品牌']})")
                    
                    # 统计设备
                    if product_id not in device_count:
                        device_count[product_id] = {
                            '设备名称': product_info['设备名称'],
                            '品牌': product_info['品牌'],
                            '主规格': product_info['主规格'],
                            '设备品类': product_info['设备品类'],
                            '单价': product_info['单价'],
                            '数量': 0
                        }
                    
                    device_count[product_id]['数量'] += 1
                    total_devices += 1
                    
                    print(f"   📊 统计设备: {product_info['设备名称']} ({product_info['品牌']})")
    
    return device_count, total_devices

def generate_product_id_report(device_count, total_devices):
    """生成产品ID识别报告"""
    
    print("\n" + "="*60)
    print("📊 产品ID识别设备统计报告")
    print("="*60)
    
    if not device_count:
        print("❌ 未通过产品ID识别到任何设备")
        print("\n💡 建议:")
        print("   1. 确保PPT中的设备使用了正确的产品ID格式")
        print("   2. 产品ID形状应该包含'_id'后缀，如'switch_1_lp_id'")
        print("   3. 或者形状名称直接是产品ID，如'switch_1_lp'")
        return
    
    # 按设备品类分组统计
    category_stats = {}
    total_cost = 0
    
    for product_id, info in device_count.items():
        category = info['设备品类']
        if category not in category_stats:
            category_stats[category] = []
        
        category_stats[category].append(info)
        total_cost += info['单价'] * info['数量']
    
    # 输出统计结果
    print(f"\n📈 总体统计:")
    print(f"   • 设备总数: {total_devices} 个")
    print(f"   • 设备种类: {len(device_count)} 种")
    print(f"   • 设备品类: {len(category_stats)} 类")
    print(f"   • 预估总价: {total_cost:.2f} 元")
    
    # 按品类输出详细信息
    for category, devices in category_stats.items():
        print(f"\n🏷️  {category}:")
        
        for device in devices:
            print(f"   📋 {device['设备名称']}")
            print(f"      • 品牌: {device['品牌']}")
            print(f"      • 规格: {device['主规格']}")
            print(f"      • 数量: {device['数量']} 个")
            print(f"      • 单价: {device['单价']} 元")
            print(f"      • 小计: {device['单价'] * device['数量']:.2f} 元")
    
    print(f"\n💰 总金额: {total_cost:.2f} 元")
    
    # 生成Excel报告
    report_data = []
    for product_id, info in device_count.items():
        report_data.append({
            '产品ID': product_id,
            '设备品类': info['设备品类'],
            '设备名称': info['设备名称'],
            '品牌': info['品牌'],
            '主规格': info['主规格'],
            '单价': info['单价'],
            '数量': info['数量'],
            '小计': info['单价'] * info['数量']
        })
    
    report_df = pd.DataFrame(report_data)
    report_path = 'E:\\Programs\\smarthome\\output\\产品ID识别报告.xlsx'
    report_df.to_excel(report_path, index=False)
    
    print(f"\n📄 详细报告已保存到: {os.path.basename(report_path)}")

def main():
    """主函数"""
    
    ppt_file_path = 'E:\\Programs\\smarthome\\output\\修复后的全屋智能方案.pptx'
    
    print("🔧 产品ID识别设备分析工具")
    print("="*60)
    print("📋 识别原理:")
    print("   • 通过形状名称中的'_id'后缀识别产品ID")
    print("   • 或者通过形状名称直接匹配产品ID")
    print("   • 与模具库中的产品ID进行精确匹配")
    print("="*60)
    
    # 通过产品ID识别设备
    device_count, total_devices = analyze_by_product_id(ppt_file_path)
    
    if device_count is not None:
        # 生成报告
        generate_product_id_report(device_count, total_devices)
    else:
        print("❌ 分析失败，请检查文件路径和格式")

if __name__ == "__main__":
    main()