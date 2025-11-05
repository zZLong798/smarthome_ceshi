#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
检查PPT中所有形状的名称和文本内容
"""

from pptx import Presentation
import re

def check_ppt_shapes():
    """检查PPT中所有形状的名称和文本内容"""
    
    ppt_path = '../全屋智能方案.pptx'
    
    try:
        prs = Presentation(ppt_path)
        print('=== PPT中所有形状名称和文本 ===')
        
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f'\n幻灯片 {slide_num}:')
            shape_count = 0
            
            for shape in slide.shapes:
                if shape.name:
                    shape_count += 1
                    print(f'  形状名称: "{shape.name}"')
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            print(f'    文本: "{text}"')
                elif shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        shape_count += 1
                        print(f'  文本形状: "{text}"')
            
            if shape_count == 0:
                print('  无形状信息')
        
        print(f'\n=== PPT文件分析完成 ===')
        
    except Exception as e:
        print(f'错误: {e}')

if __name__ == "__main__":
    check_ppt_shapes()