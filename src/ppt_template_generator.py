#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT智能家居模板生成器
创建包含预定义智能家居模具的PPT模板
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

class SmartHomeTemplateGenerator:
    """智能家居PPT模板生成器"""
    
    def __init__(self):
        self.prs = Presentation()
        self.setup_template()
    
    def setup_template(self):
        """设置模板基础样式"""
        # 设置幻灯片母版
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        
        # 添加标题幻灯片
        title_slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title_slide.shapes.title.text = "智能家居方案设计模板"
        title_slide.placeholders[1].text = "专业智能家居系统设计工具"
        
        # 添加模具库幻灯片
        self.create_shape_library_slide()
        
        # 添加设计说明幻灯片
        self.create_instruction_slide()
    
    def create_shape_library_slide(self):
        """创建智能家居模具库幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "智能家居产品模具库"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        # 创建智能开关系列
        self.create_switch_shapes(slide)
        
        # 创建传感器系列
        self.create_sensor_shapes(slide)
        
        # 创建控制器系列
        self.create_controller_shapes(slide)
    
    def create_switch_shapes(self, slide):
        """创建智能开关模具"""
        x, y = Inches(0.5), Inches(1.5)
        
        # 一键智能开关
        switch_1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1), Inches(0.6))
        switch_1.fill.solid()
        switch_1.fill.fore_color.rgb = RGBColor(74, 144, 226)  # 蓝色
        switch_1.text = "一键开关"
        switch_1.name = "smart_home_switch_1"
        
        # 二键智能开关
        switch_2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(1.5), y, Inches(1), Inches(0.6))
        switch_2.fill.solid()
        switch_2.fill.fore_color.rgb = RGBColor(74, 144, 226)
        switch_2.text = "二键开关"
        switch_2.name = "smart_home_switch_2"
        
        # 三键智能开关
        switch_3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(3), y, Inches(1), Inches(0.6))
        switch_3.fill.solid()
        switch_3.fill.fore_color.rgb = RGBColor(74, 144, 226)
        switch_3.text = "三键开关"
        switch_3.name = "smart_home_switch_3"
        
        # 四键智能开关
        switch_4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(4.5), y, Inches(1), Inches(0.6))
        switch_4.fill.solid()
        switch_4.fill.fore_color.rgb = RGBColor(74, 144, 226)
        switch_4.text = "四键开关"
        switch_4.name = "smart_home_switch_4"
    
    def create_sensor_shapes(self, slide):
        """创建传感器模具"""
        x, y = Inches(0.5), Inches(2.5)
        
        # 人体感应传感器
        sensor_1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.8), Inches(0.8))
        sensor_1.fill.solid()
        sensor_1.fill.fore_color.rgb = RGBColor(52, 168, 83)  # 绿色
        sensor_1.text = "人体感应"
        sensor_1.name = "smart_home_sensor_1"
        
        # 门窗传感器
        sensor_2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.5), y, Inches(0.8), Inches(0.8))
        sensor_2.fill.solid()
        sensor_2.fill.fore_color.rgb = RGBColor(52, 168, 83)
        sensor_2.text = "门窗感应"
        sensor_2.name = "smart_home_sensor_2"
        
        # 温湿度传感器
        sensor_3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(3), y, Inches(0.8), Inches(0.8))
        sensor_3.fill.solid()
        sensor_3.fill.fore_color.rgb = RGBColor(52, 168, 83)
        sensor_3.text = "温湿度"
        sensor_3.name = "smart_home_sensor_3"
    
    def create_controller_shapes(self, slide):
        """创建控制器模具"""
        x, y = Inches(0.5), Inches(3.7)
        
        # 智能网关
        gateway = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.2), Inches(0.8))
        gateway.fill.solid()
        gateway.fill.fore_color.rgb = RGBColor(255, 153, 0)  # 橙色
        gateway.text = "智能网关"
        gateway.name = "smart_home_gateway"
        
        # 场景控制器
        controller = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.8), y, Inches(1.2), Inches(0.8))
        controller.fill.solid()
        controller.fill.fore_color.rgb = RGBColor(255, 153, 0)
        controller.text = "场景控制"
        controller.name = "smart_home_controller"
    
    def create_instruction_slide(self):
        """创建使用说明幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        slide.shapes.title.text = "使用说明"
        
        content = slide.placeholders[1]
        content.text = """
1. 模具库使用：
   • 复制模具到设计页面
   • 系统自动识别产品类型
   
2. 设计规范：
   • 每个房间单独设计
   • 标注产品数量和位置
   
3. 报价生成：
   • 使用PPT到Excel工具
   • 自动计算总价和利润
        """
    
    def save_template(self, filename="smart_home_template.pptx"):
        """保存模板文件"""
        self.prs.save(filename)
        print(f"✅ 智能家居PPT模板已生成: {filename}")
        print("📋 模板包含:")
        print("   • 标题幻灯片")
        print("   • 智能家居模具库")
        print("   • 使用说明")
        print("   • 7种预定义产品模具")
        return filename

def main():
    """主函数"""
    print("=== PPT智能家居模板生成器 ===")
    print("正在创建专业智能家居设计模板...")
    
    generator = SmartHomeTemplateGenerator()
    template_file = generator.save_template()
    
    print("\n🎯 模板特性:")
    print("   • 预定义智能家居产品模具")
    print("   • 标准化的设计规范")
    print("   • 支持自动识别和报价")
    print("   • 专业的外观设计")
    
    print(f"\n📁 模板文件: {template_file}")
    print("💡 使用方法: 打开模板文件，复制模具到设计页面进行方案设计")

if __name__ == "__main__":
    main()