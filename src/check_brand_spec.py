#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
检查智能开关的品牌和规格信息
"""

from ppt_to_excel_bridge import PPTtoExcelBridge
from pptx import Presentation

def check_brand_spec_details():
    """检查智能开关的详细品牌规格信息"""
    print("=== 智能开关品牌规格详细信息 ===")
    print()
    
    # 1. 查看产品库中的详细信息
    bridge = PPTtoExcelBridge()
    
    print("📋 产品库中的智能开关信息:")
    print("=" * 50)
    
    switch_products = []
    for product_id, info in bridge.product_library.items():
        if 'switch' in product_id:
            switch_products.append((product_id, info))
    
    # 按产品ID排序显示
    switch_products.sort(key=lambda x: x[0])
    
    for product_id, info in switch_products:
        print(f"🔹 产品ID: {product_id}")
        print(f"   名称: {info.get('name', '未知')}")
        print(f"   品牌: {info.get('brand', '未知')}")
        print(f"   型号: {info.get('model', '未知')}")
        print(f"   规格: {info.get('spec', '未知')}")
        print(f"   价格: {info.get('price', 0)}元")
        print()
    
    print("=" * 50)
    
    # 2. 检查您方案中实际使用的开关信息
    print("\n🔍 您方案中使用的智能开关:")
    print("=" * 50)
    
    prs = Presentation('../全屋智能方案.pptx')
    slide = prs.slides[6]  # 第7张幻灯片
    
    switch_count = {}
    
    for shape in slide.shapes:
        if shape.name and 'switch' in shape.name.lower():
            # 提取产品ID
            name = shape.name.lower()
            product_id = None
            
            if 'smart_home_switch_1' in name:
                product_id = 'switch_1'
            elif 'smart_home_switch_2' in name:
                product_id = 'switch_2'
            elif 'smart_home_switch_3' in name:
                product_id = 'switch_3'
            elif 'smart_home_switch_4' in name:
                product_id = 'switch_4'
            
            if product_id:
                if product_id not in switch_count:
                    switch_count[product_id] = 0
                switch_count[product_id] += 1
    
    # 显示您方案中使用的开关
    if switch_count:
        for product_id, count in switch_count.items():
            info = bridge.product_library.get(product_id, {})
            print(f"🔸 {info.get('name', '未知')} ({product_id}): {count}个")
            print(f"   品牌: {info.get('brand', '未知')}")
            print(f"   型号: {info.get('model', '未知')}")
            print(f"   规格: {info.get('spec', '未知')}")
            print(f"   单价: {info.get('price', 0)}元")
            print()
    else:
        print("未找到智能开关设备")
    
    print("=" * 50)
    
    # 3. 生成采购清单格式的信息
    print("\n📊 采购清单格式信息:")
    print("=" * 50)
    
    total_price = 0
    for product_id, count in switch_count.items():
        info = bridge.product_library.get(product_id, {})
        price = info.get('price', 0)
        subtotal = count * price
        total_price += subtotal
        
        print(f"{info.get('name', '未知')}:")
        print(f"  品牌: {info.get('brand', '未知')}")
        print(f"  型号: {info.get('model', '未知')}")
        print(f"  规格: {info.get('spec', '未知')}")
        print(f"  数量: {count}个")
        print(f"  单价: {price}元")
        print(f"  小计: {subtotal}元")
        print()
    
    print(f"💎 总计: {sum(switch_count.values())} 个设备")
    print(f"💰 总价: {total_price} 元")

if __name__ == "__main__":
    check_brand_spec_details()