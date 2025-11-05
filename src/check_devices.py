#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
准确统计智能家居设备数量和类型
"""

from pptx import Presentation

def count_smart_devices(ppt_path):
    """准确统计PPT中的智能设备"""
    print("=== 智能设备详细统计 ===")
    
    prs = Presentation(ppt_path)
    
    # 产品映射表
    product_map = {
        'switch_1': '一键智能开关',
        'switch_2': '二键智能开关', 
        'switch_3': '三键智能开关',
        'switch_4': '四键智能开关',
        'sensor_motion': '人体感应传感器',
        'sensor_door': '门窗传感器',
        'sensor_temp': '温湿度传感器',
        'gateway': '智能网关'
    }
    
    price_map = {
        'switch_1': 79,
        'switch_2': 89,
        'switch_3': 99,
        'switch_4': 109,
        'sensor_motion': 65,
        'sensor_door': 45,
        'sensor_temp': 55,
        'gateway': 299
    }
    
    device_count = {}
    slide_devices = {}
    
    # 扫描所有幻灯片
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_found = False
        
        for shape in slide.shapes:
            if shape.name:
                name = shape.name.lower()
                
                # 检查是否为智能设备
                product_id = None
                
                # 格式1: smart_home_switch_1
                if 'smart_home' in name:
                    parts = name.split('smart_home_')[1].split('_')
                    if len(parts) >= 2:
                        product_id = f"{parts[0]}_{parts[1]}"
                
                # 格式2: switch_1
                elif 'switch' in name:
                    for i in range(1, 5):
                        if str(i) in name:
                            product_id = f"switch_{i}"
                            break
                
                # 格式3: 传感器和网关
                elif 'sensor' in name or '传感器' in name:
                    if 'motion' in name or '人体' in name:
                        product_id = 'sensor_motion'
                    elif 'door' in name or '门窗' in name:
                        product_id = 'sensor_door'
                    elif 'temp' in name or '温湿' in name:
                        product_id = 'sensor_temp'
                
                elif 'gateway' in name or '网关' in name:
                    product_id = 'gateway'
                
                if product_id and product_id in product_map:
                    # 统计设备
                    if product_id not in device_count:
                        device_count[product_id] = 0
                    device_count[product_id] += 1
                    
                    # 记录幻灯片位置
                    if slide_num not in slide_devices:
                        slide_devices[slide_num] = []
                    slide_devices[slide_num].append(product_id)
                    
                    slide_found = True
                    
                    print(f"  幻灯片{slide_num}: 发现 {product_map[product_id]}")
    
    # 显示统计结果
    print("\n📊 设备统计结果:")
    print("=" * 40)
    
    total_count = 0
    total_price = 0
    
    if device_count:
        for product_id, count in device_count.items():
            name = product_map[product_id]
            price = price_map[product_id]
            subtotal = count * price
            
            print(f"  {name} ({product_id}): {count}个")
            print(f"    单价: {price}元, 小计: {subtotal}元")
            
            total_count += count
            total_price += subtotal
    else:
        print("  未发现智能设备")
    
    print("=" * 40)
    print(f"  总计: {total_count} 个设备")
    print(f"  总价: {total_price} 元")
    
    # 显示设备分布
    if slide_devices:
        print("\n📍 设备分布位置:")
        for slide_num, devices in slide_devices.items():
            device_names = [product_map[pid] for pid in devices]
            print(f"  幻灯片{slide_num}: {', '.join(device_names)}")
    
    return device_count

if __name__ == "__main__":
    count_smart_devices('../全屋智能方案.pptx')