#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Excel到PPT模具库转换器
根据Excel表格自动生成智能家居模具库PPT
"""

import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
import requests
from PIL import Image
import io
from urllib.parse import urlparse

class ExcelToPPTConverter:
    """Excel到PPT模具库转换器"""
    
    def __init__(self, image_folder="assets/images"):
        """
        初始化转换器
        
        Args:
            image_folder: 图片存储文件夹
        """
        self.image_folder = image_folder
        self.ensure_image_folder()
        
        # 设置项目根目录
        import os
        self.project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        
        # 模具布局配置
        self.layout_config = {
            'slides_per_page': 6,  # 每页幻灯片显示的模具数量
            'mold_width': Inches(2.5),
            'mold_height': Inches(2.0),
            'margin': Inches(0.5),
            'spacing': Inches(0.3)
        }
        
        # 颜色配置
        self.color_scheme = {
            'title': RGBColor(0, 51, 102),      # 深蓝色
            'subtitle': RGBColor(102, 102, 102), # 灰色
            'price': RGBColor(204, 0, 0),       # 红色
            'brand': RGBColor(0, 102, 51),      # 绿色
            'background': RGBColor(255, 255, 255) # 白色
        }
    
    def ensure_image_folder(self):
        """确保图片文件夹存在"""
        if not os.path.exists(self.image_folder):
            os.makedirs(self.image_folder)
            print(f"创建图片文件夹: {self.image_folder}")
    
    def parse_excel_image_mapping(self, excel_file_path):
        """
        解析Excel文件中的图片映射关系
        
        Args:
            excel_file_path: Excel文件路径
            
        Returns:
            dict: 图片ID到映射信息的字典
        """
        try:
            # 导入我们之前创建的图片映射解析脚本的功能
            import sys
            import os
            
            # 添加项目根目录到Python路径
            project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            if project_root not in sys.path:
                sys.path.insert(0, project_root)
            
            # 检查是否存在图片映射解析脚本
            mapping_script_path = os.path.join(project_root, 'parse_correct_mapping_final.py')
            if not os.path.exists(mapping_script_path):
                print("图片映射解析脚本不存在，跳过映射解析")
                return None
            
            # 动态导入映射解析功能
            import importlib.util
            spec = importlib.util.spec_from_file_location("image_mapping", mapping_script_path)
            mapping_module = importlib.util.module_from_spec(spec)
            
            # 执行映射解析脚本
            spec.loader.exec_module(mapping_module)
            
            # 调用映射解析功能
            mapping_result = mapping_module.parse_excel_image_mapping(excel_file_path)
            
            if mapping_result:
                print(f"成功解析图片映射，找到 {len(mapping_result)} 个映射关系")
                return mapping_result
            else:
                print("图片映射解析返回空结果")
                return None
                
        except Exception as e:
            print(f"图片映射解析失败: {e}")
            return None
    
    def extract_embedded_images(self, sheet, temp_dir="temp_excel_images"):
        """
        提取Excel中嵌入的图片并保存为临时文件
        
        Args:
            sheet: Excel工作表对象
            temp_dir: 临时文件目录
            
        Returns:
            dict: 行号到图片路径的映射
        """
        image_mapping = {}
        
        # 首先尝试使用图片映射解析功能
        excel_file_path = sheet.parent.path if hasattr(sheet, 'parent') and hasattr(sheet.parent, 'path') else None
        if excel_file_path and os.path.exists(excel_file_path):
            print("尝试使用图片映射解析功能...")
            mapping_result = self.parse_excel_image_mapping(excel_file_path)
            if mapping_result:
                print(f"图片映射解析成功，找到 {len(mapping_result)} 个图片映射")
                # 将映射结果转换为行号到图片路径的映射
                for image_id, mapping_info in mapping_result.items():
                    row_num = mapping_info.get('row_number')
                    if row_num and mapping_info.get('actual_file_path'):
                        image_mapping[row_num] = mapping_info['actual_file_path']
                        print(f"行{row_num} -> {mapping_info['actual_file_path']}")
                return image_mapping
        
        # 如果图片映射解析失败，回退到原始方法
        print("图片映射解析失败，使用原始嵌入图片提取方法...")
        
        # 创建临时目录
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
        
        # 提取嵌入的图片
        for i, img in enumerate(sheet._images):
            try:
                # 获取图片锚点位置
                if hasattr(img.anchor, '_from'):
                    cell_ref = img.anchor._from
                    row_num = cell_ref.row
                    col_num = cell_ref.col
                    
                    print(f"发现嵌入图片: 行{row_num}, 列{col_num}")
                    
                    # 保存图片到临时文件
                    image_filename = f"excel_image_row{row_num}_col{col_num}.png"
                    image_path = os.path.join(temp_dir, image_filename)
                    
                    # 保存图片
                    with open(image_path, 'wb') as f:
                        f.write(img._data())
                    
                    # 映射到行号
                    image_mapping[row_num] = image_path
                    print(f"图片已保存: {image_path}")
                    
            except Exception as e:
                print(f"提取图片失败: {e}")
        
        return image_mapping
    
    def extract_image_urls_from_formulas(self, sheet, excel_file_path=None):
        """
        从DISPIMG公式中提取图片URL或路径信息
        
        Args:
            sheet: Excel工作表对象
            excel_file_path: Excel文件路径（可选）
            
        Returns:
            dict: 行号到图片URL/路径的映射
        """
        image_url_mapping = {}
        
        # 首先尝试使用图片映射解析功能
        if not excel_file_path:
            # 尝试从sheet对象获取文件路径（可能不可行）
            excel_file_path = sheet.parent.path if hasattr(sheet, 'parent') and hasattr(sheet.parent, 'path') else None
        
        if excel_file_path and os.path.exists(excel_file_path):
            print("尝试使用图片映射解析功能...")
            mapping_result = self.parse_excel_image_mapping(excel_file_path)
            if mapping_result:
                print(f"图片映射解析成功，找到 {len(mapping_result)} 个图片映射")
                # 将映射结果转换为行号到图片路径的映射
                for image_name, mapping_info in mapping_result.items():
                    row_num = mapping_info.get('row_number')
                    if row_num:
                        # 优先查找重命名的设备图片
                        renamed_image_path = self.find_renamed_device_image(row_num)
                        if renamed_image_path and os.path.exists(renamed_image_path):
                            image_url_mapping[row_num] = {
                                'type': 'mapped_image',
                                'image_name': image_name,
                                'actual_file_path': renamed_image_path,
                                'file_name': mapping_info['file_name'],
                                'description': mapping_info['description']
                            }
                            print(f"行{row_num} -> 重命名设备图片: {renamed_image_path}")
                        elif mapping_info.get('actual_file_path'):
                            # 如果重命名图片不存在，使用映射解析的临时路径
                            image_url_mapping[row_num] = {
                                'type': 'mapped_image',
                                'image_name': image_name,
                                'actual_file_path': mapping_info['actual_file_path'],
                                'file_name': mapping_info['file_name'],
                                'description': mapping_info['description']
                            }
                            print(f"行{row_num} -> 映射图片: {mapping_info['actual_file_path']}")
                return image_url_mapping
        
        # 如果图片映射解析失败，使用原始方法
        print("图片映射解析失败，使用原始公式提取方法...")
        
        # 获取表头
        headers = []
        for col in range(1, sheet.max_column + 1):
            header = sheet.cell(row=1, column=col).value
            headers.append(header if header else f"列{col}")
        
        # 找到图片列
        image_col_index = None
        for i, header in enumerate(headers):
            if '图片' in str(header):
                image_col_index = i + 1
                print(f"找到图片列: {header} (列{image_col_index})")
                break
        
        if not image_col_index:
            print("未找到图片列")
            return image_url_mapping
        
        # 扫描所有行，提取图片URL
        for row in range(2, sheet.max_row + 1):
            try:
                image_cell = sheet.cell(row=row, column=image_col_index)
                
                # 检查单元格是否有值
                if not image_cell.value:
                    continue
                
                # 调试信息：显示单元格内容和数据类型
                print(f"行{row}: 单元格值={repr(image_cell.value)}, 数据类型={image_cell.data_type}")
                    
                # 如果是公式，尝试解析DISPIMG
                if image_cell.data_type == 'f' and 'DISPIMG' in str(image_cell.value):
                    formula = image_cell.value
                    print(f"行{row}: 发现DISPIMG公式: {formula}")
                    
                    # 对于DISPIMG公式，我们无法直接获取图片，但可以记录行号用于后续处理
                    # 这里我们标记该行有图片需求，但需要其他方式获取图片
                    image_url_mapping[row] = {
                        'type': 'dispimg_formula',
                        'formula': formula,
                        'has_image_need': True
                    }
                    
                # 如果是文本，可能是图片URL或路径
                elif image_cell.data_type == 's':
                    url_or_path = image_cell.value.strip()
                    if url_or_path and url_or_path != 'None':
                        print(f"行{row}: 发现图片URL/路径: {url_or_path}")
                        image_url_mapping[row] = {
                            'type': 'url_or_path',
                            'value': url_or_path
                        }
                        
            except Exception as e:
                print(f"处理行{row}的图片信息失败: {e}")
        
        print(f"从公式中提取到 {len(image_url_mapping)} 个图片信息")
        return image_url_mapping
    
    def extract_image_from_other_columns(self, product, headers, row_number=None):
        """
        从产品数据的其他列中提取图片URL或路径，或提供默认图片
        
        Args:
            product: 产品数据字典
            headers: 表头列表
            row_number: Excel中的行号
            
        Returns:
            str: 图片URL或路径，如果未找到返回默认图片路径
        """
        device_name = product.get('设备名称', 'unknown')
        print(f"extract_image_from_other_columns被调用: 行号={row_number}, 设备名称={device_name}")
        
        # 1. 首先检查是否有DISPIMG公式，尝试解析图片ID
        possible_image_columns = [
            '图片URL', '图片路径', '图片地址', 'image_url', 'image_path',
            '产品图片', '设备图片', '图片', 'photo', 'image'
        ]
        
        for col_name in possible_image_columns:
            if col_name in product and product[col_name]:
                value = product[col_name]
                if isinstance(value, str) and value.strip() and value.strip() != 'None':
                    print(f"检查列 '{col_name}': 值={repr(value)}")
                    # 如果是DISPIMG公式，尝试解析图片ID
                    if 'DISPIMG' in value:
                        print(f"检测到DISPIMG公式: {value}")
                        # 提取图片ID - 匹配格式：=_xlfn.DISPIMG("ID_...",1)
                        import re
                        pattern = r'DISPIMG\("([^"]+)",1\)'
                        match = re.search(pattern, value)
                        if match:
                            image_id = match.group(1)
                            print(f"提取的图片ID: {image_id}")
                            # 尝试根据图片ID和设备名称查找本地图片
                            local_image = self.find_local_image_by_id_and_name(image_id, device_name, row_number)
                            if local_image:
                                print(f"找到本地图片: {local_image}")
                                return local_image
                            else:
                                print(f"未找到图片ID {image_id} 对应的本地图片")
                                # 如果找不到对应的Excel图片，使用默认图片
                                return self.get_default_image_path({'设备品类': product.get('设备品类', 'unknown'), '设备名称': device_name})
                        else:
                            print(f"DISPIMG公式模式不匹配: {value}")
                        continue
                    
                    # 检查是否是URL
                    import re
                    url_pattern = re.compile(r'^https?://')
                    if url_pattern.match(value.strip()):
                        return value.strip()
                    
                    # 检查是否是相对路径或绝对路径
                    if os.path.exists(value.strip()):
                        return value.strip()
                    
                    # 检查是否是相对路径（相对于项目根目录）
                    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
                    relative_path = os.path.join(project_root, value.strip())
                    if os.path.exists(relative_path):
                        return relative_path
        
        # 2. 尝试从其他文本列中查找URL模式
        for header, value in product.items():
            if isinstance(value, str) and value.strip():
                # 查找URL模式
                import re
                url_pattern = re.compile(r'https?://[^\s]+')
                matches = url_pattern.findall(value)
                if matches:
                    # 检查是否是图片URL（常见图片扩展名）
                    image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']
                    for url in matches:
                        if any(url.lower().endswith(ext) for ext in image_extensions):
                            return url
        
        # 3. 最后，为每个设备创建独特的默认图片，而不是使用品类默认图片
        return self.create_unique_device_image(product)
    
    def find_local_image_by_id_and_name(self, image_id, device_name, row_number):
        """
        根据图片ID、设备名称和行号查找本地图片文件
        
        Args:
            image_id: 图片ID（来自DISPIMG公式）
            device_name: 设备名称
            row_number: Excel中的行号
            
        Returns:
            str: 本地图片路径，如果未找到返回None
        """
        # 首先检查temp_excel_images目录中的Excel提取图片
        temp_excel_dir = 'temp_excel_images'
        if os.path.exists(temp_excel_dir) and row_number is not None:
            # Excel数据行号从2开始，但图片文件行号从1开始，需要转换
            # 例如：Excel第2行数据对应图片文件中的row1
            image_row_number = row_number - 1
            
            # 查找对应行号的图片文件
            possible_patterns = [
                f"excel_image_row{image_row_number}_col11.png",  # 第11列是设备图片列
                f"excel_image_row{image_row_number}_col*.png",   # 任何列
                f"excel_image_row{image_row_number}_*.png"       # 任何列和格式
            ]
            
            print(f"查找Excel图片: 行号={row_number}, 转换后图片行号={image_row_number}")
            print(f"查找模式: {possible_patterns}")
            
            import glob
            for pattern in possible_patterns:
                search_pattern = os.path.join(temp_excel_dir, pattern)
                matching_files = glob.glob(search_pattern)
                print(f"搜索模式: {search_pattern}, 找到文件: {matching_files}")
                if matching_files:
                    # 返回第一个匹配的文件
                    file_path = matching_files[0]
                    print(f"找到Excel提取的设备图片: {file_path}")
                    return file_path
            
            print(f"未找到行{row_number}对应的Excel图片文件")
            # 如果找不到Excel图片，直接返回None，让上层逻辑处理默认图片
            return None
        else:
            print(f"temp_excel_images目录不存在或行号为空: {row_number}")
            return None
    
    def find_renamed_device_image(self, row_num):
        """根据行号查找重命名的设备图片"""
        renamed_dir = os.path.join(self.project_root, 'renamed_device_images')
        if not os.path.exists(renamed_dir):
            return None
        
        # 行号从1开始，但Excel数据行通常从第2行开始（第1行是标题）
        # 所以行号需要加1来对应PDID
        pdid = row_num + 1
        
        # 查找以pdid开头的文件
        for filename in os.listdir(renamed_dir):
            if filename.startswith(f"pdid{pdid}_"):
                file_path = os.path.join(renamed_dir, filename)
                if os.path.exists(file_path):
                    return file_path
        
        return None
    
    def find_device_image(self, product):
        """
        查找设备图片，优先使用映射关系文件进行查找
        
        Args:
            product: 产品数据字典
            
        Returns:
            str: 设备图片路径，如果找不到返回None
        """
        # 获取设备信息
        device_name = product.get('设备名称', '')
        product_id = product.get('产品ID', '')
        
        # 如果没有设备信息，返回None
        if not device_name or not product_id:
            return None
        
        # 获取项目根目录下的images目录（与Excel文件同目录）
        project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        images_dir = os.path.join(project_root, 'images')
        
        # 如果images目录不存在，返回None
        if not os.path.exists(images_dir):
            return None
        
        # 首先尝试使用映射关系文件查找
        mapping_file_path = os.path.join(images_dir, 'image_mapping.json')
        if os.path.exists(mapping_file_path):
            try:
                # 导入映射关系生成器
                from image_mapping_generator import ImageMappingGenerator
                
                # 创建映射生成器实例
                excel_file_path = os.path.join(project_root, '智能家居模具库.xlsx')
                mapping_generator = ImageMappingGenerator(excel_file_path, images_dir)
                
                # 根据PDID查找图片路径
                image_path = mapping_generator.get_image_by_pdid(product_id)
                if image_path and os.path.exists(image_path):
                    print(f"通过映射关系找到设备图片: {os.path.basename(image_path)}")
                    return image_path
                
            except Exception as e:
                print(f"映射关系查找失败，使用备用查找方式: {e}")
        
        # 备用方案：尝试多种文件名匹配模式
        device_short_name = product.get('设备简称', '')
        possible_filenames = []
        
        # 模式1: pdid{id}_{设备简称拼音}
        if device_short_name:
            # 简单的拼音转换（将中文转换为拼音首字母或拼音）
            pinyin_name = self.convert_to_pinyin(device_short_name)
            possible_filenames.append(f"pdid{product_id}_{pinyin_name}.png")
        
        # 模式2: pdid{id}_{设备名称拼音}
        pinyin_full_name = self.convert_to_pinyin(device_name)
        possible_filenames.append(f"pdid{product_id}_{pinyin_full_name}.png")
        
        # 模式3: 仅pdid{id}
        possible_filenames.append(f"pdid{product_id}.png")
        
        # 模式4: 设备名称相关文件名
        safe_name = "".join(c for c in device_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
        safe_name = safe_name.replace(' ', '_')
        possible_filenames.append(f"{safe_name}_{product_id}.png")
        
        # 在images目录中查找匹配的文件
        for filename in possible_filenames:
            image_path = os.path.join(images_dir, filename)
            if os.path.exists(image_path):
                print(f"通过文件名匹配找到设备图片: {filename}")
                return image_path
        
        # 如果没有找到匹配的图片，返回None
        return None
    
    def convert_to_pinyin(self, chinese_text):
        """
        将中文文本转换为拼音（简化版本）
        
        Args:
            chinese_text: 中文文本
            
        Returns:
            str: 拼音字符串
        """
        if not chinese_text:
            return ""
        
        # 简单的拼音映射表（仅包含常见汉字）
        pinyin_map = {
            '一': 'yi', '二': 'er', '三': 'san', '四': 'si', '五': 'wu',
            '六': 'liu', '七': 'qi', '八': 'ba', '九': 'jiu', '十': 'shi',
            '键': 'jian', '开': 'kai', '关': 'guan', '智': 'zhi', '能': 'neng',
            '灯': 'deng', '具': 'ju', '窗': 'chuang', '帘': 'lian', '传': 'chuan',
            '感': 'gan', '器': 'qi', '家': 'jia', '电': 'dian', '门': 'men',
            '锁': 'suo', '传': 'chuan', '感': 'gan', '器': 'qi'
        }
        
        # 将中文转换为拼音
        pinyin_result = ""
        for char in chinese_text:
            if char in pinyin_map:
                pinyin_result += pinyin_map[char]
            elif char.isalnum():
                pinyin_result += char
            elif char in (' ', '-', '_'):
                pinyin_result += char
        
        # 如果转换失败，使用设备名称的简化版本
        if not pinyin_result:
            safe_name = "".join(c for c in chinese_text if c.isalnum() or c in (' ', '-', '_')).rstrip()
            pinyin_result = safe_name.replace(' ', '_')
        
        return pinyin_result
    
    def create_unique_device_image(self, product):
        """
        为每个设备创建独特的默认图片，而不是使用品类默认图片
        
        Args:
            product: 产品数据字典
            
        Returns:
            str: 独特的设备图片路径
        """
        # 创建设备图片目录
        device_images_dir = os.path.join(self.project_root, 'images')
        if not os.path.exists(device_images_dir):
            os.makedirs(device_images_dir)
        
        # 获取设备信息
        device_type = product.get('设备品类', 'unknown')
        device_name = product.get('设备名称', 'unknown')
        product_id = product.get('产品ID', 'unknown')
        
        # 生成独特的文件名
        safe_name = "".join(c for c in device_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
        safe_name = safe_name.replace(' ', '_')
        unique_filename = f"{safe_name}_{product_id}.png"
        unique_image_path = os.path.join(device_images_dir, unique_filename)
        
        # 如果独特的图片已存在，直接返回
        if os.path.exists(unique_image_path):
            return unique_image_path
        
        # 创建独特的设备图片
        return self.create_default_image(unique_image_path, device_type, device_name)
    
    def get_default_image_path(self, product):
        """
        为没有图片的设备提供默认图片路径
        
        Args:
            product: 产品数据字典
            
        Returns:
            str: 默认图片路径
        """
        # 创建默认图片目录
        default_images_dir = os.path.join(self.project_root, 'default_images')
        if not os.path.exists(default_images_dir):
            os.makedirs(default_images_dir)
        
        # 根据设备类型选择不同的默认图片
        device_type = product.get('设备品类', 'unknown')
        device_name = product.get('设备名称', 'unknown')
        
        # 设备类型到默认图片的映射
        default_images = {
            '智能开关': 'smart_switch.png',
            '智能灯具': 'smart_light.png', 
            '全屋WiFi': 'wifi_router.png',
            '智能门锁': 'smart_lock.png',
            '智能窗帘': 'smart_curtain.png',
            '智能传感器': 'sensor.png',
            '智能家电': 'smart_appliance.png',
            '其他': 'default_device.png'
        }
        
        # 获取对应的默认图片文件名
        image_filename = default_images.get(device_type, 'default_device.png')
        default_image_path = os.path.join(default_images_dir, image_filename)
        
        # 如果默认图片不存在，创建一个简单的占位图片
        if not os.path.exists(default_image_path):
            self.create_default_image(default_image_path, device_type, device_name)
        
        print(f"为设备 {device_name} 使用默认图片: {default_image_path}")
        return default_image_path
    
    def create_default_image(self, image_path, device_type, device_name):
        """
        创建默认占位图片
        
        Args:
            image_path: 图片保存路径
            device_type: 设备类型
            device_name: 设备名称
        """
        try:
            from PIL import Image, ImageDraw, ImageFont
            
            # 创建更大的图片（400x300像素），确保有足够的数据
            img = Image.new('RGB', (400, 300), color=(255, 255, 255))  # 白色背景
            draw = ImageDraw.Draw(img)
            
            # 绘制边框
            draw.rectangle([(10, 10), (390, 290)], outline=(200, 200, 200), width=2)
            
            # 绘制设备图标（简单的几何图形）
            center_x, center_y = 200, 120
            
            # 根据设备类型绘制不同的图标
            if '开关' in device_type:
                # 开关图标：矩形
                draw.rectangle([center_x-30, center_y-20, center_x+30, center_y+20], 
                              fill=(100, 150, 255), outline=(50, 100, 200), width=2)
            elif '灯' in device_type or '照明' in device_type:
                # 灯具图标：圆形
                draw.ellipse([center_x-25, center_y-25, center_x+25, center_y+25], 
                            fill=(255, 200, 100), outline=(200, 150, 50), width=2)
            elif '窗帘' in device_type:
                # 窗帘图标：波浪线
                for i in range(5):
                    x = center_x - 40 + i * 20
                    draw.arc([x, center_y-15, x+15, center_y+15], 0, 180, fill=(150, 200, 150), width=3)
            elif '路由' in device_type or 'WiFi' in device_type:
                # 路由器图标：信号波
                for i in range(4):
                    radius = 15 + i * 8
                    draw.arc([center_x-radius, center_y-radius, center_x+radius, center_y+radius], 
                            0, 180, fill=(100, 200, 100), width=2)
            else:
                # 默认图标：齿轮
                draw.ellipse([center_x-20, center_y-20, center_x+20, center_y+20], 
                           fill=(200, 200, 200), outline=(150, 150, 150), width=2)
            
            # 尝试加载字体
            try:
                # 尝试多种字体
                fonts_to_try = ['arial.ttf', 'arialbd.ttf', 'times.ttf', 'calibri.ttf']
                font = None
                for font_name in fonts_to_try:
                    try:
                        font = ImageFont.truetype(font_name, 16)
                        break
                    except:
                        continue
                
                if font is None:
                    font = ImageFont.load_default()
            except:
                font = ImageFont.load_default()
            
            # 绘制设备类型（居中）
            text_width = draw.textlength(device_type, font=font)
            draw.text((center_x - text_width/2, center_y + 50), device_type, 
                     fill=(0, 0, 0), font=font)
            
            # 绘制设备名称（截断过长的名称）
            short_name = device_name[:12] + '...' if len(device_name) > 12 else device_name
            text_width = draw.textlength(short_name, font=font)
            draw.text((center_x - text_width/2, center_y + 80), short_name, 
                     fill=(100, 100, 100), font=font)
            
            # 保存为高质量PNG图片
            img.save(image_path, 'PNG', quality=95)
            
            # 验证图片文件大小
            file_size = os.path.getsize(image_path)
            print(f"创建默认图片: {image_path} (大小: {file_size} 字节)")
            
            if file_size < 1000:
                print(f"警告: 图片文件过小，可能有问题")
                
        except ImportError:
            # 如果PIL不可用，复制现有的默认图片
            default_image = os.path.join(self.project_root, 'default_images', 'default_device.png')
            if os.path.exists(default_image):
                import shutil
                shutil.copy2(default_image, image_path)
                print(f"复制默认图片: {image_path}")
            else:
                # 创建简单的占位文件
                with open(image_path, 'wb') as f:
                    # 创建一个小的有效PNG文件
                    f.write(b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x00^\xdd\x86\x1f\x00\x00\x00\x00IEND\xaeB`\x82')
                print(f"创建PNG占位文件: {image_path}")
        except Exception as e:
            print(f"创建默认图片失败: {e}")
            # 确保至少有一个有效的文件
            if not os.path.exists(image_path):
                with open(image_path, 'wb') as f:
                    f.write(b'')
    
    def process_product_image(self, product, row_number=None):
        """
        处理产品图片，按照新流程：
        1. 检查renamed_device_images目录中是否存在pdid{产品ID}_设备简称.png格式的图片
        2. 如果存在，直接使用该图片
        3. 如果不存在，则调用image_mapping_generator.py中的generate_image_mapping方法
        4. 使用生成的映射关系，将图片复制到renamed_device_images目录并重命名
        5. 最后使用重命名后的图片
        
        Args:
            product: 产品数据字典
            row_number: Excel中的行号
            
        Returns:
            str: 图片路径，如果处理失败返回None
        """
        # 获取产品信息 - 适配多种可能的列名
        product_id = product.get('产品ID') or product.get('型号') or str(product.get('序号', ''))
        short_name = product.get('设备简称') or product.get('设备') or ''
        device_name = product.get('设备名称') or product.get('设备') or ''
        
        # 如果没有标准产品ID，尝试从设备名称推断
        if not product_id or product_id in ['【白色 四开】智能开关-AG玻璃-Mesh2.0零火版', '灰色']:
            if '二键' in device_name or '二键' in short_name:
                product_id = '6'  # 二键开关对应pdid6
            elif '四键' in device_name or '四键' in short_name:
                product_id = '8'  # 四键开关对应pdid8
            else:
                print(f"无法从设备名称推断产品ID: {device_name}")
                return None
        
        if not product_id:
            print(f"产品ID为空，无法处理图片: {device_name}")
            return None
        
        # 1. 检查renamed_device_images目录中是否存在pdid{产品ID}_设备简称.png格式的图片
        renamed_dir = os.path.join(self.project_root, 'renamed_device_images')
        if not os.path.exists(renamed_dir):
            os.makedirs(renamed_dir)
            print(f"创建重命名图片目录: {renamed_dir}")
        
        # 生成预期的重命名文件名
        safe_short_name = "".join(c for c in short_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
        safe_short_name = safe_short_name.replace(' ', '_')
        expected_filename = f"pdid{product_id}_{safe_short_name}.png"
        expected_path = os.path.join(renamed_dir, expected_filename)
        
        # 2. 如果存在，直接使用该图片
        if os.path.exists(expected_path):
            print(f"找到已重命名的图片: {expected_path}")
            return expected_path
        
        # 3. 如果不存在，则调用image_mapping_generator.py中的generate_image_mapping方法
        try:
            # 导入image_mapping_generator模块
            import sys
            sys.path.insert(0, self.project_root)
            from image_mapping_generator import ImageMappingGenerator
            
            # 创建ImageMappingGenerator实例
            # 首先尝试在项目根目录查找，如果不存在则在backup_main_flow目录查找
            excel_path = os.path.join(os.path.dirname(self.project_root), '智能家居模具库.xlsx')
            if not os.path.exists(excel_path):
                excel_path = os.path.join(self.project_root, '智能家居模具库.xlsx')
            
            images_dir = os.path.join(os.path.dirname(self.project_root), 'images')
            if not os.path.exists(images_dir):
                images_dir = os.path.join(self.project_root, 'images')
            
            if not os.path.exists(excel_path):
                print(f"Excel文件不存在: {excel_path}")
                return None
                
            if not os.path.exists(images_dir):
                print(f"图片目录不存在: {images_dir}")
                return None
            
            mapping_generator = ImageMappingGenerator(excel_path, images_dir)
            
            # 获取产品ID对应的图片路径
            print(f"查找产品ID {product_id} 的图片映射...")
            source_path = mapping_generator.get_image_by_pdid(str(product_id))
            
            if not source_path:
                print(f"未找到产品ID {product_id} 的图片映射")
                return None
            if not os.path.exists(source_path):
                print(f"源图片文件不存在: {source_path}")
                return None
            
            # 复制并重命名图片
            import shutil
            shutil.copy2(source_path, expected_path)
            print(f"图片已复制并重命名: {source_path} -> {expected_path}")
            
            # 5. 返回重命名后的图片路径
            return expected_path
            
        except ImportError:
            print("无法导入image_mapping_generator模块")
            return None
        except Exception as e:
            print(f"处理图片映射时出错: {e}")
            return None
    
    def read_excel_data(self, excel_path):
        """
        读取Excel数据
        
        Args:
            excel_path: Excel文件路径
            
        Returns:
            list: 产品数据列表
        """
        try:
            workbook = openpyxl.load_workbook(excel_path)
            sheet = workbook.active
            
            # 读取表头
            headers = []
            for col in range(1, sheet.max_column + 1):
                header = sheet.cell(row=1, column=col).value
                headers.append(header if header else f"列{col}")
            
            # 读取数据
            products = []
            for row in range(2, sheet.max_row + 1):
                product = {}
                valid_row = False
                
                for col, header in enumerate(headers, 1):
                    value = sheet.cell(row=row, column=col).value
                    product[header] = value
                    
                    # 检查是否有效行（至少有一个非空值）
                    if value and not valid_row:
                        valid_row = True
                
                if valid_row:
                    # 检查是否启用
                    is_enabled = product.get('是否启用', True)
                    if is_enabled in [True, '是', '启用', '1', 1]:
                        # 不再处理图片，将在generate_ppt_from_excel中使用新的图片处理流程
                        products.append(product)
            
            print(f"从Excel读取到 {len(products)} 个启用的产品")
            return products
            
        except Exception as e:
            print(f"读取Excel文件失败: {e}")
            return []
    
    def download_image(self, image_url, product_name):
        """
        处理产品图片
        
        Args:
            image_url: 图片URL或本地路径
            product_name: 产品名称（用于文件名）
            
        Returns:
            str: 本地图片路径，如果处理失败返回None
        """
        if not image_url:
            return None
        
        try:
            # 检查是否是本地文件路径
            if os.path.exists(image_url):
                print(f"使用本地图片: {image_url}")
                return image_url
            
            # 检查是否是相对路径（相对于项目根目录）
            project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            relative_path = os.path.join(project_root, image_url)
            if os.path.exists(relative_path):
                print(f"使用相对路径图片: {relative_path}")
                return relative_path
            
            # 如果是URL，下载图片
            parsed_url = urlparse(image_url)
            if parsed_url.scheme in ['http', 'https']:
                # 解析URL获取文件扩展名
                file_ext = os.path.splitext(parsed_url.path)[1] or '.jpg'
                
                # 生成安全的文件名
                safe_name = "".join(c for c in product_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
                safe_name = safe_name.replace(' ', '_')
                
                # 本地文件路径
                local_path = os.path.join(self.image_folder, f"{safe_name}{file_ext}")
                
                # 如果文件已存在，直接返回路径
                if os.path.exists(local_path):
                    return local_path
                
                # 下载图片
                response = requests.get(image_url, timeout=10)
                response.raise_for_status()
                
                # 保存图片
                with open(local_path, 'wb') as f:
                    f.write(response.content)
                
                print(f"下载图片成功: {local_path}")
                return local_path
            else:
                print(f"无法识别的图片路径: {image_url}")
                return None
            
        except Exception as e:
            print(f"处理图片失败 {image_url}: {e}")
            return None
    
    def create_mold_shape(self, slide, product, position_x, position_y):
        """
        创建模具形状
        
        Args:
            slide: PPT幻灯片对象
            product: 产品数据
            position_x: X坐标
            position_y: Y坐标
            
        Returns:
            shape: 创建的形状对象
        """
        width = self.layout_config['mold_width']
        height = self.layout_config['mold_height']
        
        # 创建模具容器（矩形）
        mold_box = slide.shapes.add_shape(
            1,  # 矩形
            position_x, position_y, width, height
        )
        
        # 设置模具名称（智能标记）
        product_type = product.get('设备品类', 'unknown')
        product_name = product.get('设备简称', product.get('设备名称', 'unknown'))
        mold_box.name = f"smart_home_{product_type}_{product_name}"
        
        # 设置边框样式
        mold_box.line.color.rgb = RGBColor(200, 200, 200)
        mold_box.line.width = Pt(1)
        
        # 添加产品图片
        image_path = product.get('local_image_path')
        if image_path and os.path.exists(image_path):
            try:
                # 图片位置（顶部）
                img_width = width - Inches(0.4)
                img_height = height * 0.5
                img_x = position_x + (width - img_width) / 2
                img_y = position_y + Inches(0.2)
                
                slide.shapes.add_picture(image_path, img_x, img_y, img_width, img_height)
            except Exception as e:
                print(f"添加图片失败 {image_path}: {e}")
        
        # 添加产品信息文本
        self.add_product_info(slide, product, position_x, position_y, width, height)
        
        return mold_box
    
    def add_product_info(self, slide, product, x, y, width, height):
        """添加产品信息文本"""
        
        # 产品名称
        product_name = product.get('设备名称', '未知产品')
        name_textbox = slide.shapes.add_textbox(
            x + Inches(0.1), y + height * 0.5,
            width - Inches(0.2), Inches(0.4)
        )
        name_frame = name_textbox.text_frame
        name_frame.text = product_name
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.size = Pt(12)
        name_frame.paragraphs[0].font.color.rgb = self.color_scheme['title']
        name_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 品牌和规格
        brand = product.get('品牌', '')
        spec = product.get('主规格', '')
        spec_text = f"{brand} {spec}".strip()
        
        if spec_text:
            spec_textbox = slide.shapes.add_textbox(
                x + Inches(0.1), y + height * 0.5 + Inches(0.4),
                width - Inches(0.2), Inches(0.3)
            )
            spec_frame = spec_textbox.text_frame
            spec_frame.text = spec_text
            spec_frame.paragraphs[0].font.size = Pt(10)
            spec_frame.paragraphs[0].font.color.rgb = self.color_scheme['brand']
            spec_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 价格
        price = product.get('单价', 0)
        if price:
            price_textbox = slide.shapes.add_textbox(
                x + Inches(0.1), y + height - Inches(0.5),
                width - Inches(0.2), Inches(0.3)
            )
            price_frame = price_textbox.text_frame
            price_frame.text = f"¥{price}"
            price_frame.paragraphs[0].font.bold = True
            price_frame.paragraphs[0].font.size = Pt(14)
            price_frame.paragraphs[0].font.color.rgb = self.color_scheme['price']
            price_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def create_slide_title(self, slide, title):
        """创建幻灯片标题"""
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.size = Pt(16)  # 减小字体大小，避免过大
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.color_scheme['title']
        
        # 调整标题位置和大小，避免跑到外面去
        title_shape.left = Inches(0.5)  # 左对齐，留出边距
        title_shape.top = Inches(0.1)    # 上移标题，更紧凑
        title_shape.width = Inches(9)    # 设置合适宽度
        title_shape.height = Inches(0.6)  # 减小高度，更紧凑
    
    def group_products_by_category(self, products):
        """按设备品类分组产品"""
        categories = {}
        
        for product in products:
            category = product.get('设备品类', '其他')
            if category not in categories:
                categories[category] = []
            categories[category].append(product)
        
        return categories
    
    def generate_ppt_from_excel(self, excel_path, ppt_path=None):
        """
        从Excel生成PPT模具库
        
        Args:
            excel_path: Excel文件路径
            ppt_path: 输出的PPT文件路径
            
        Returns:
            bool: 是否成功生成
        """
        if ppt_path is None:
            base_name = os.path.splitext(excel_path)[0]
            ppt_path = f"{base_name}_模具库.pptx"
        
        # 1. 读取Excel数据
        products = self.read_excel_data(excel_path)
        if not products:
            print("没有找到有效的产品数据")
            return False
        
        # 2. 处理产品图片（使用新的图片处理流程）
        print("处理产品图片...")
        for i, product in enumerate(products):
            # 使用新的图片处理流程
            image_path = self.process_product_image(product, i+2)  # Excel数据从第2行开始
            if image_path:
                product['local_image_path'] = image_path
                device_name = product.get('设备名称') or product.get('设备') or 'unknown'
                print(f"为产品 {device_name} 设置图片: {image_path}")
            else:
                device_name = product.get('设备名称') or product.get('设备') or 'unknown'
                print(f"产品 {device_name} 图片处理失败")
        
        # 3. 按品牌/品类分组
        products_by_brand = {}
        for product in products:
            brand = product.get('品牌', '其他')
            category = product.get('设备品类', '其他')
            
            # 智能开关按品牌分类，其他设备按品类分类
            if category == '智能开关':
                if brand not in products_by_brand:
                    products_by_brand[brand] = []
                products_by_brand[brand].append(product)
            else:
                if category not in products_by_brand:
                    products_by_brand[category] = []
                products_by_brand[category].append(product)
        
        print(f"按品牌/品类分组: {list(products_by_brand.keys())}")
        
        # 4. 创建PPT
        prs = Presentation()
        
        # 封面页
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        self.create_slide_title(slide, "智能家居模具库")
        
        subtitle = slide.placeholders[1]
        subtitle.text = f"基于Excel自动生成 - 共{len(products)}个产品"
        
        # 为每个品牌/品类创建幻灯片
        for group_name, group_products in products_by_brand.items():
            category = group_products[0].get('设备品类', '其他')
            if category == '智能开关':
                # 智能开关使用品牌名称
                slide_title = f"{group_name}智能开关系列"
            else:
                # 其他设备使用品类名称
                slide_title = f"{category}系列"
            
            print(f"生成 {slide_title} 幻灯片...")
            
            # 创建品类标题页
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            self.create_slide_title(slide, slide_title)
            
            # 计算布局
            products_per_slide = self.layout_config['slides_per_page']
            
            for i in range(0, len(group_products), products_per_slide):
                batch_products = group_products[i:i + products_per_slide]
                
                # 创建产品展示页
                if i > 0:  # 第一页已经创建了品类标题页
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页
                
                # 添加产品模具
                self.add_products_to_slide(slide, batch_products, group_name)
        
        # 5. 保存PPT
        try:
            prs.save(ppt_path)
            print(f"PPT模具库生成成功: {ppt_path}")
            print(f"总计生成 {len(prs.slides)} 张幻灯片")
            return True
        except Exception as e:
            print(f"保存PPT失败: {e}")
            return False
    
    def add_smart_mark_to_shape(self, shape, product):
        """为形状添加智能标记，用于统计数量"""
        try:
            # 设置智能标记属性
            product_type = product.get('设备品类', 'unknown')
            short_name = product.get('设备简称', product.get('设备名称', 'unknown'))
            brand = product.get('品牌', 'unknown')
            
            # 生成与识别工具兼容的智能标记
            # 识别工具支持格式：smart_home_switch_2 或 switch_2
            if "开关" in product_type or "switch" in product_type.lower():
                # 从设备简称中提取开关键数
                if "一键" in short_name or "1键" in short_name:
                    shape.name = "smart_home_switch_1"
                elif "二键" in short_name or "2键" in short_name:
                    shape.name = "smart_home_switch_2"
                elif "三键" in short_name or "3键" in short_name:
                    shape.name = "smart_home_switch_3"
                elif "四键" in short_name or "4键" in short_name:
                    shape.name = "smart_home_switch_4"
                else:
                    # 默认使用简化格式
                    shape.name = "switch_1"
            else:
                # 其他产品类型使用简化格式
                shape.name = f"smart_home_{product_type}"
            
            # 添加自定义属性（如果支持）
            try:
                # 尝试设置自定义属性用于统计
                shape._element.set('smart_home_product', 'true')
                shape._element.set('product_type', product_type)
                shape._element.set('short_name', short_name)
                shape._element.set('brand', brand)
            except:
                # 如果自定义属性不支持，使用名称标记即可
                pass
                
        except Exception as e:
            print(f"添加智能标记失败: {e}")
    
    def add_products_to_slide(self, slide, products, group_name):
        """在幻灯片上添加产品模具"""
        # 计算布局参数
        products_per_row = 4  # 增加每行显示数量
        products_per_col = 2
        
        # 计算每个产品的位置和大小
        slide_width = Inches(10)
        slide_height = Inches(7.5)
        margin = Inches(0.3)  # 减小边距
        
        content_width = slide_width - 2 * margin
        content_height = slide_height - 2 * margin
        
        cell_width = content_width / products_per_row
        cell_height = content_height / products_per_col
        
        # 图片和简称组成一个组，其他信息单独显示
        group_height_ratio = 0.6  # 调整图片+简称组的高度比例
        info_height_ratio = 0.4   # 增加其他信息的高度比例
        
        group_height = cell_height * group_height_ratio
        info_height = cell_height * info_height_ratio
        
        # 添加每个产品
        for idx, product in enumerate(products):
            row = idx // products_per_row
            col = idx % products_per_row
            
            # 计算位置
            left = margin + col * cell_width
            top = margin + row * cell_height
            
            # 创建图片+简称组（用于复制和统计）
            group_left = left + Inches(0.05)
            group_top = top + Inches(0.05)
            group_width = cell_width - Inches(0.1)
            
            # 添加产品图片（固定大小0.9cm x 0.9cm）
            picture = None
            if product.get('local_image_path'):
                try:
                    print(f"尝试添加图片: {product['local_image_path']}")
                    
                    # 检查图片文件是否存在
                    import os
                    if not os.path.exists(product['local_image_path']):
                        print(f"图片文件不存在: {product['local_image_path']}")
                    else:
                        print(f"图片文件存在，大小: {os.path.getsize(product['local_image_path'])} 字节")
                    
                    # 固定图片大小：0.9cm x 0.9cm
                    img_width = Inches(0.9 / 2.54)  # 厘米转英寸
                    img_height = Inches(0.9 / 2.54)
                    
                    # 计算图片位置（居中）
                    img_left = group_left + (group_width - img_width) / 2
                    img_top = group_top + Inches(0.08)  # 增大顶部边距，增加与文字的距离
                    
                    # 添加图片
                    picture = slide.shapes.add_picture(
                        product['local_image_path'],
                        img_left, img_top, img_width, img_height
                    )
                    device_name = product.get('设备名称') or product.get('设备') or 'unknown'
                    print(f"图片添加成功: {device_name}, 图片类型: {type(picture)}, 是否有image属性: {hasattr(picture, 'image')}")
                    
                    # 检查图片形状的属性
                    if hasattr(picture, 'image'):
                        print(f"图片数据大小: {len(picture.image.blob) if hasattr(picture.image, 'blob') else 'N/A'}")
                    else:
                        print("图片形状没有image属性")
                    
                    # 添加边框（实线，颜色RGB(255, 217, 102)）
                    picture.line.color.rgb = RGBColor(255, 217, 102)
                    picture.line.width = Pt(2)  # 边框宽度
                    
                    # 添加智能标记到图片上（用于统计数量）
                    self.add_smart_mark_to_shape(picture, product)
                    
                except Exception as e:
                    device_name = product.get('设备名称') or product.get('设备') or 'unknown'
                    print(f"添加图片失败 {device_name}: {e}")
                    picture = None
            
            # 添加设备简称文本（带橙色填充背景）
            short_name = product.get('设备简称') or product.get('设备') or product.get('设备名称') or '未知'
            
            # 如果没有图片，使用默认位置
            if picture:
                # 计算文字背景框位置（增大与图片的距离）
                text_bg_left = img_left - Inches(0.02)
                text_bg_top = img_top + img_height + Inches(0.05)  # 增大与图片的距离，避免重叠
                text_bg_width = img_width + Inches(0.04)
            else:
                # 没有图片时的默认位置
                text_bg_left = group_left
                text_bg_top = group_top + Inches(0.08)
                text_bg_width = group_width
            
            text_bg_height = Inches(0.2)  # 减小文字区域高度
            
            # 添加橙色背景填充
            text_bg_shape = slide.shapes.add_shape(
                1,  # 矩形
                text_bg_left, text_bg_top, text_bg_width, text_bg_height
            )
            text_bg_shape.fill.solid()
            text_bg_shape.fill.fore_color.rgb = RGBColor(197, 90, 17)  # 橙色填充
            text_bg_shape.line.fill.background()  # 无边框
            
            # 添加设备简称文本框（在背景框内部）
            short_name_left = text_bg_left + Inches(0.01)
            short_name_top = text_bg_top + Inches(0.01)
            short_name_width = text_bg_width - Inches(0.02)
            short_name_height = text_bg_height - Inches(0.02)
            
            short_name_textbox = slide.shapes.add_textbox(
                short_name_left, short_name_top, short_name_width, short_name_height
            )
            text_frame = short_name_textbox.text_frame
            text_frame.clear()
            
            # 设置文本框垂直对齐为居中
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 添加设备简称
            p = text_frame.paragraphs[0]
            p.text = short_name
            p.font.size = Pt(6)  # 字号设为6
            p.font.color.rgb = RGBColor(255, 217, 102)  # 颜色改为RGB(255, 217, 102)
            p.alignment = PP_ALIGN.CENTER
            
            # 为简称文本添加智能标记
            self.add_smart_mark_to_shape(short_name_textbox, product)
            
            # 添加pdid标签（透明背景，无边框，最小号字体）
            pdid = product.get('产品ID') or product.get('型号') or str(idx + 1)  # 使用产品ID或型号或默认编号
            
            # 计算pdid标签位置（在设备简称下方）
            pdid_left = text_bg_left
            pdid_top = text_bg_top + text_bg_height + Inches(0.02)  # 在简称下方
            pdid_width = text_bg_width
            pdid_height = Inches(0.15)  # 较小高度
            
            # 创建pdid标签文本框
            pdid_textbox = slide.shapes.add_textbox(pdid_left, pdid_top, pdid_width, pdid_height)
            pdid_textbox.name = f"pdid_label_{pdid}"
            
            # 设置pdid文本框样式
            pdid_frame = pdid_textbox.text_frame
            pdid_frame.clear()
            
            # 设置文本框垂直对齐为居中
            pdid_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 添加pdid文本
            p_pdid = pdid_frame.paragraphs[0]
            p_pdid.text = f"pdid: {pdid}"
            p_pdid.font.size = Pt(6)  # 最小号字体
            p_pdid.font.color.rgb = RGBColor(255, 255, 255)  # 白色文字
            p_pdid.font.fill.background()  # 文本填充设为透明
            p_pdid.alignment = PP_ALIGN.CENTER
            
            # 设置pdid文本框背景 - 透明不填充
            pdid_fill = pdid_textbox.fill
            pdid_fill.background()  # 透明背景
            
            # 设置pdid文本框边框 - 无边框
            pdid_line = pdid_textbox.line
            pdid_line.fill.background()  # 透明边框
            
            # 为pdid标签添加智能标记
            self.add_smart_mark_to_shape(pdid_textbox, product)
            
            # 将设备图片、橙色背景、简称文本和pdid标签组合在一起
            try:
                # 收集所有需要组合的形状
                group_shapes = []
                if picture:
                    group_shapes.append(picture)
                group_shapes.extend([text_bg_shape, short_name_textbox, pdid_textbox])
                
                # 计算组的边界框
                shape_lefts = [shape.left for shape in group_shapes]
                shape_tops = [shape.top for shape in group_shapes]
                shape_rights = [shape.left + shape.width for shape in group_shapes]
                shape_bottoms = [shape.top + shape.height for shape in group_shapes]
                
                group_left = min(shape_lefts)
                group_top = min(shape_tops)
                group_right = max(shape_rights)
                group_bottom = max(shape_bottoms)
                group_width = group_right - group_left
                group_height = group_bottom - group_top
                
                # 创建组
                device_group = slide.shapes.add_group_shape(group_shapes)
                
                # 设置组的位置和大小
                device_group.left = group_left
                device_group.top = group_top
                device_group.width = group_width
                device_group.height = group_height
                
                # 为整个设备组添加智能标记
                self.add_smart_mark_to_shape(device_group, product)
                
                device_name = product.get('设备名称') or product.get('设备') or 'unknown'
                print(f"设备组组合成功: {device_name}, 包含{len(group_shapes)}个形状")
                
            except Exception as e:
                device_name = product.get('设备名称') or product.get('设备') or 'unknown'
                print(f"组合设备组失败 {device_name}: {e}")
                # 如果组合失败，仍然为各个形状添加智能标记
                if picture:
                    self.add_smart_mark_to_shape(picture, product)
                self.add_smart_mark_to_shape(text_bg_shape, product)
                self.add_smart_mark_to_shape(short_name_textbox, product)
                self.add_smart_mark_to_shape(pdid_textbox, product)
            
            # 添加其他信息（设备名称、规格、价格）
            info_left = left + Inches(0.05)
            info_top = top + group_height + Inches(0.15)  # 大幅增加间距，方便点击设备组
            info_width = cell_width - Inches(0.1)
            info_height_adjusted = info_height - Inches(0.1)  # 调整高度
            
            info_textbox = slide.shapes.add_textbox(info_left, info_top, info_width, info_height_adjusted)
            info_frame = info_textbox.text_frame
            info_frame.word_wrap = True
            info_frame.vertical_anchor = MSO_ANCHOR.TOP  # 顶部对齐
            
            # 添加设备名称（小字体）
            device_name = product.get('设备名称') or product.get('设备') or ''
            if device_name and device_name != short_name:
                p = info_frame.paragraphs[0]
                p.text = device_name
                p.font.size = Pt(7)  # 减小字体
                p.font.color.rgb = RGBColor(128, 128, 128)  # 灰色
                p.space_after = Pt(2)  # 增加段落间距
            
            # 添加规格信息
            if product.get('主规格'):
                p = info_frame.add_paragraph()
                p.text = f"规格: {product['主规格']}"
                p.font.size = Pt(7)
                p.font.color.rgb = RGBColor(128, 128, 128)
                p.space_after = Pt(2)
            
            # 添加价格信息
            if product.get('单价'):
                p = info_frame.add_paragraph()
                p.text = f"价格: ¥{product['单价']}"
                p.font.size = Pt(7)
                p.font.color.rgb = RGBColor(255, 0, 0)  # 红色价格
                p.space_after = Pt(2)
    
    def create_sample_excel(self, excel_path):
        """创建示例Excel文件"""
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "智能家居产品库"
        
        # 表头
        headers = [
            "设备品类", "设备名称", "设备简称", "是否启用", 
            "单价", "品牌", "主规格", "单位", 
            "渠道", "采购链接", "设备图片"
        ]
        
        for col, header in enumerate(headers, 1):
            sheet.cell(row=1, column=col, value=header)
        
        # 示例数据
        sample_data = [
            ["智能开关", "一键智能开关", "一键开关", "是", 79, "颜工", "86型", "个", "电商", "https://example.com/switch1", "https://example.com/switch1.jpg"],
            ["智能开关", "二键智能开关", "二键开关", "是", 89, "颜工", "86型", "个", "电商", "https://example.com/switch2", "https://example.com/switch2.jpg"],
            ["传感器", "人体感应传感器", "人体感应", "是", 65, "颜工", "Zigbee", "个", "电商", "https://example.com/motion", "https://example.com/motion.jpg"],
            ["传感器", "门窗传感器", "门窗感应", "是", 45, "颜工", "Zigbee", "个", "电商", "https://example.com/door", "https://example.com/door.jpg"],
            ["控制器", "智能网关", "智能网关", "是", 299, "颜工", "WiFi+Zigbee", "台", "电商", "https://example.com/gateway", "https://example.com/gateway.jpg"]
        ]
        
        for row, data in enumerate(sample_data, 2):
            for col, value in enumerate(data, 1):
                sheet.cell(row=row, column=col, value=value)
        
        # 保存文件
        workbook.save(excel_path)
        print(f"示例Excel文件已创建: {excel_path}")

# 使用示例
if __name__ == "__main__":
    import sys
    
    # 处理命令行参数
    excel_path = "智能家居模具库.xlsx"  # 默认值
    ppt_path = "智能家居模具库_修复测试2.pptx"  # 默认值
    
    # 解析命令行参数
    if len(sys.argv) >= 2:
        for i, arg in enumerate(sys.argv):
            if arg == "--excel" and i + 1 < len(sys.argv):
                excel_path = sys.argv[i + 1]
            elif arg == "--output" and i + 1 < len(sys.argv):
                ppt_path = sys.argv[i + 1]
    
    converter = ExcelToPPTConverter()
    
    print(f"使用Excel文件: {excel_path}")
    print(f"输出PPT文件: {ppt_path}")
    
    success = converter.generate_ppt_from_excel(excel_path, ppt_path)
    
    if success:
        print(f"PPT模具库生成成功: {ppt_path}")
        print("使用用户文档生成完成！")
    else:
        print("PPT生成失败")