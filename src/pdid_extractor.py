#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pdid标签提取模块 - 专门处理pdid: 1格式的标签提取
"""

from pptx import Presentation
import re
from typing import List, Dict, Optional


class PDIDExtractor:
    """pdid标签提取器"""
    
    def __init__(self, ppt_path: str):
        """
        初始化pdid标签提取器
        
        Args:
            ppt_path: PPT文件路径
        """
        self.ppt_path = ppt_path
        self.presentation = None
        
    def load_presentation(self) -> bool:
        """
        加载PPT文件
        
        Returns:
            bool: 是否成功加载
        """
        try:
            self.presentation = Presentation(self.ppt_path)
            print(f"✅ 成功加载PPT文件: {self.ppt_path}")
            print(f"📊 幻灯片数量: {len(self.presentation.slides)}")
            return True
        except Exception as e:
            print(f"❌ 加载PPT文件失败: {e}")
            return False
    
    def extract_pdid_labels(self) -> Dict[int, List[Dict]]:
        """
        提取PPT中的pdid标签（按照项目规则：从组内精准匹配）
        
        Returns:
            Dict[int, List[Dict]]: 幻灯片索引到pdid标签信息的映射
        """
        if self.presentation is None:
            print("❌ 请先加载PPT文件")
            return {}
        
        pdid_labels = {}
        
        print("\n🔍 开始按照项目规则提取PPT中的pdid标签...")
        print("📋 规则：从组内精准匹配pdid标签，不依赖形状名称")
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            print(f"\n📄 扫描第{slide_idx + 1}张幻灯片:")
            
            slide_labels = []
            
            # 检查所有形状，包括嵌套的组合形状
            for shape in slide.shapes:
                shape_name = shape.name if hasattr(shape, 'name') else ""
                
                # 检查形状的文本内容
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    try:
                        text = shape.text.strip() if shape.text else ""
                        if text:
                            # 精准匹配pdid格式：pdid: 数字
                            pdid_values = self._match_pdid_format(text)
                            if pdid_values:
                                for pdid_value in pdid_values:
                                    label_info = {
                                        'shape': shape,
                                        'name': shape_name,
                                        'text': text,
                                        'pdid': pdid_value,
                                        'type': 'direct_text',
                                        'position': {
                                            'left': shape.left if hasattr(shape, 'left') else 0,
                                            'top': shape.top if hasattr(shape, 'top') else 0,
                                            'width': shape.width if hasattr(shape, 'width') else 0,
                                            'height': shape.height if hasattr(shape, 'height') else 0
                                        }
                                    }
                                    slide_labels.append(label_info)
                                    print(f"   ✅ 发现pdid标签: {text} (形状: {shape_name})")
                    except Exception as e:
                        print(f"   ⚠️ 处理形状文本失败: {e}")
                
                # 检查形状是否是组合形状，并递归检查子形状
                if hasattr(shape, 'shapes'):
                    # 标记为组合形状
                    print(f"   🔍 检查组合形状: {shape_name}")
                    
                    for sub_shape in shape.shapes:
                        sub_name = sub_shape.name if hasattr(sub_shape, 'name') else ""
                        
                        # 检查子形状的文本内容
                        if hasattr(sub_shape, 'has_text_frame') and sub_shape.has_text_frame:
                            try:
                                text = sub_shape.text.strip() if sub_shape.text else ""
                                if text:
                                    # 精准匹配pdid格式：pdid: 数字
                                    pdid_values = self._match_pdid_format(text)
                                    if pdid_values:
                                        for pdid_value in pdid_values:
                                            label_info = {
                                                'shape': sub_shape,
                                                'name': sub_name,
                                                'text': text,
                                                'pdid': pdid_value,
                                                'type': 'group_text',
                                                'parent_group': shape_name,
                                                'position': {
                                                    'left': sub_shape.left if hasattr(sub_shape, 'left') else 0,
                                                    'top': sub_shape.top if hasattr(sub_shape, 'top') else 0,
                                                    'width': sub_shape.width if hasattr(sub_shape, 'width') else 0,
                                                    'height': sub_shape.height if hasattr(sub_shape, 'height') else 0
                                                }
                                            }
                                            slide_labels.append(label_info)
                                            print(f"   ✅ 从组内发现pdid标签: {text} (组合: {shape_name}, 子形状: {sub_name})")
                            except Exception as e:
                                print(f"   ⚠️ 处理组合子形状文本失败: {e}")
            
            pdid_labels[slide_idx] = slide_labels
            print(f"   📊 本页发现pdid标签: {len(slide_labels)}个")
        
        total_labels = sum(len(labels) for labels in pdid_labels.values())
        print(f"\n📈 总计发现pdid标签: {total_labels}个")
        
        if total_labels == 0:
            print("⚠️ 未发现任何pdid标签，请检查PPT文件中的pdid格式是否正确")
            
        return pdid_labels
    
    def _extract_pdid_from_slide(self, slide, slide_idx: int) -> List[Dict]:
        """
        从单张幻灯片中提取pdid标签
        
        Args:
            slide: 幻灯片对象
            slide_idx: 幻灯片索引
            
        Returns:
            List[Dict]: pdid标签信息列表
        """
        slide_labels = []
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            try:
                text = shape.text.strip()
                if not text:
                    continue
                    
                # 使用正则表达式匹配pdid: 1格式（冒号后有一个空格）
                pdid_matches = self._match_pdid_format(text)
                
                if pdid_matches:
                    for pdid_value in pdid_matches:
                        label_info = {
                            'slide_index': slide_idx,
                            'slide_number': slide_idx + 1,
                            'shape_name': shape.name,
                            'text_content': text,
                            'pdid': pdid_value,
                            'position': {
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height
                            }
                        }
                        slide_labels.append(label_info)
                        print(f"   ✅ 发现pdid标签: pdid: {pdid_value} (形状: {shape.name})")
                        
            except Exception as e:
                print(f"   ⚠️ 处理形状 {shape.name} 时出错: {e}")
                continue
        
        return slide_labels
    
    def _match_pdid_format(self, text: str) -> List[int]:
        """
        匹配pdid标签格式（按照项目规则：pdid: 数字）
        
        Args:
            text: 文本内容
            
        Returns:
            List[int]: 匹配到的pdid值列表，未匹配到返回空列表
        """
        pdid_values = []
        
        # 严格按照项目规则匹配：pdid: 数字（冒号后有一个空格）
        pdid_pattern = r'pdid:\s*(\d+)'
        match = re.search(pdid_pattern, text, re.IGNORECASE)
        if match:
            pdid_values.append(int(match.group(1)))
            print(f"      🔍 匹配到标准pdid格式: {match.group(0)}")
        
        # 如果标准格式未匹配，尝试宽松匹配（允许空格变化）
        if not pdid_values:
            pdid_pattern_loose = r'pdid\s*:\s*(\d+)'
            match_loose = re.search(pdid_pattern_loose, text, re.IGNORECASE)
            if match_loose:
                pdid_values.append(int(match_loose.group(1)))
                print(f"      🔍 匹配到宽松pdid格式: {match_loose.group(0)}")
        
        return pdid_values
    
    def _create_product_id_mapping(self) -> Dict[str, int]:
        """
        创建产品ID映射关系
        
        Returns:
            Dict[str, int]: 字符串产品ID到数字产品ID的映射
        """
        # 根据之前分析的结果创建映射
        mapping = {
            'switch_1_yl': 1,
            'switch_2_yl': 2, 
            'switch_3_yl': 3,
            'switch_4_yl': 4,
            'switch_1': 1,
            'switch_2': 2,
            'switch_3': 3,
            'switch_4': 4,
            '_id_1': 1,
            '_id_2': 2,
            '_id_3': 3,
            '_id_4': 4,
        }
        return mapping
    
    def _extract_product_id_from_shape_name(self, shape_name: str) -> Optional[int]:
        """
        从形状名称中提取产品ID（已弃用，按照项目规则应从组内文本匹配）
        
        Args:
            shape_name: 形状名称
            
        Returns:
            Optional[int]: 始终返回None，因为项目规则要求从组内文本匹配
        """
        # 按照项目规则，不从形状名称中提取PDID
        # PDID应该从组内文本内容中精准匹配
        return None
    
    def get_pdid_list(self) -> List[int]:
        """
        获取所有提取的pdid值列表
        
        Returns:
            List[int]: pdid值列表
        """
        pdid_labels = self.extract_pdid_labels()
        pdid_values = []
        
        # 从字典中提取所有pdid值
        for slide_labels in pdid_labels.values():
            for label in slide_labels:
                pdid_values.append(label['pdid'])
        
        # 去重并排序
        return sorted(list(set(pdid_values)))
    
    def save_pdid_report(self, output_path: str = "pdid_extraction_report.json") -> bool:
        """
        保存pdid提取报告
        
        Args:
            output_path: 输出文件路径
            
        Returns:
            bool: 是否成功保存
        """
        try:
            pdid_labels = self.extract_pdid_labels()
            
            # 计算总标签数和唯一pdid值
            total_labels = 0
            all_pdid_values = []
            
            # 创建可序列化的报告数据
            serializable_labels = {}
            
            for slide_idx, slide_labels in pdid_labels.items():
                serializable_labels[slide_idx] = []
                
                for label in slide_labels:
                    total_labels += 1
                    all_pdid_values.append(label['pdid'])
                    
                    # 创建可序列化的标签信息（不包含Shape对象）
                    serializable_label = {
                        'name': label['name'],
                        'text': label['text'],
                        'pdid': label['pdid'],
                        'type': label['type'],
                        'position': label['position']
                    }
                    serializable_labels[slide_idx].append(serializable_label)
            
            unique_pdid_values = list(set(all_pdid_values))
            
            report = {
                'ppt_file': self.ppt_path,
                'total_pdid_labels': total_labels,
                'unique_pdid_values': unique_pdid_values,
                'pdid_labels': serializable_labels
            }
            
            import json
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(report, f, ensure_ascii=False, indent=2)
            
            print(f"\n💾 pdid提取报告已保存至: {output_path}")
            return True
            
        except Exception as e:
            print(f"❌ 保存pdid提取报告失败: {e}")
            return False


def test_pdid_extractor(ppt_path=None):
    """测试pdid标签提取器"""
    if ppt_path is None:
        ppt_path = "../全屋智能方案.pptx"
    
    extractor = PDIDExtractor(ppt_path)
    
    if extractor.load_presentation():
        pdid_labels = extractor.extract_pdid_labels()
        pdid_list = extractor.get_pdid_list()
        
        print(f"\n📋 提取的pdid值列表: {pdid_list}")
        
        # 保存报告
        extractor.save_pdid_report("pdid_extraction_report.json")
        
        return pdid_list
    
    return []


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        # 使用命令行参数指定的PPT文件
        ppt_path = sys.argv[1]
        test_pdid_extractor(ppt_path)
    else:
        # 使用默认PPT文件
        test_pdid_extractor()